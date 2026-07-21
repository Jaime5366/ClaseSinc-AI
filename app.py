import streamlit as st
import os
import tempfile
import re
import time
import json
from pathlib import Path
from pptx import Presentation
from pypdf import PdfReader
try:
    from moviepy.editor import VideoFileClip
except ImportError:
    from moviepy import VideoFileClip
from dotenv import load_dotenv
from google import genai
from google.genai import errors
from fpdf import FPDF
from openai import OpenAI

# Cargar variables de entorno del archivo .env si existe
load_dotenv()

# Configuración de página de Streamlit
st.set_page_config(
    page_title="ClaseSinc AI - Resúmenes Inteligentes de Clases",
    page_icon="🎓",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Estilo CSS personalizado para lograr una interfaz premium y moderna
st.markdown("""
<style>
    /* Estilos globales */
    .main {
        background-color: #0f111a;
        color: #e2e8f0;
    }
    
    /* Degradado premium para el título principal */
    .title-gradient {
        background: linear-gradient(135deg, #6366f1 0%, #a855f7 50%, #ec4899 100%);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
        font-family: 'Outfit', 'Inter', sans-serif;
        font-weight: 800;
        font-size: 3rem;
        text-align: center;
        margin-bottom: 0.2rem;
    }
    
    .subtitle-text {
        text-align: center;
        color: #94a3b8;
        font-size: 1.15rem;
        margin-bottom: 2rem;
    }
    
    /* Tarjetas personalizadas */
    .css-card {
        background-color: #1e293b;
        border-radius: 12px;
        padding: 24px;
        border: 1px solid #334155;
        box-shadow: 0 4px 6px -1px rgba(0, 0, 0, 0.1), 0 2px 4px -1px rgba(0, 0, 0, 0.06);
        margin-bottom: 1.5rem;
    }
</style>
""", unsafe_allow_html=True)

# --- Base de Datos Local ---
DB_DIR = Path("saved_classes")
DB_DIR.mkdir(exist_ok=True)

READINGS_DIR = Path("saved_readings")
READINGS_DIR.mkdir(exist_ok=True)

def migrate_legacy_files():
    general_classes = DB_DIR / "General"
    general_readings = READINGS_DIR / "General"
    general_classes.mkdir(parents=True, exist_ok=True)
    general_readings.mkdir(parents=True, exist_ok=True)
    
    # Migrar clases
    for p in DB_DIR.glob("*.json"):
        if p.is_file():
            try:
                dest = general_classes / p.name
                if not dest.exists():
                    p.rename(dest)
                else:
                    p.unlink()
            except Exception:
                pass
                
    # Migrar lecturas
    for p in READINGS_DIR.glob("*.json"):
        if p.is_file():
            try:
                dest = general_readings / p.name
                if not dest.exists():
                    p.rename(dest)
                else:
                    p.unlink()
            except Exception:
                pass

# Ejecutar migración automática de archivos heredados al iniciar
migrate_legacy_files()

def clean_subject_name(name):
    return re.sub(r'[\\/*?:"<>|]', "", name).strip()

def list_subjects():
    subjects = set()
    
    # 1. Intentar consultar desde Supabase DB
    try:
        from supabase_client import get_supabase_client
        client = get_supabase_client()
        res = client.table("documents").select("title").execute()
        if res.data:
            for row in res.data:
                t = row.get("title", "")
                if t.startswith("[") and "]" in t:
                    subj = t[1:t.index("]")].strip()
                    if subj:
                        subjects.add(subj)
    except Exception as e:
        print(f"[Supabase list_subjects Notice] {e}")

    # 2. Consultar disco local (fallback / dev)
    if DB_DIR.exists():
        for p in DB_DIR.iterdir():
            if p.is_dir():
                subjects.add(p.name)
    if READINGS_DIR.exists():
        for p in READINGS_DIR.iterdir():
            if p.is_dir():
                subjects.add(p.name)
                
    if not subjects:
        subjects.add("General")
    return sorted(list(subjects))

def create_subject(name):
    cleaned = clean_subject_name(name)
    if cleaned:
        (DB_DIR / cleaned).mkdir(parents=True, exist_ok=True)
        (READINGS_DIR / cleaned).mkdir(parents=True, exist_ok=True)
        return cleaned
    return None

def rename_subject(old_name, new_name):
    cleaned_new = clean_subject_name(new_name)
    if not cleaned_new or cleaned_new == old_name:
        return False
        
    old_classes = DB_DIR / old_name
    new_classes = DB_DIR / cleaned_new
    if old_classes.exists() and not new_classes.exists():
        old_classes.rename(new_classes)
    else:
        new_classes.mkdir(parents=True, exist_ok=True)
        
    old_readings = READINGS_DIR / old_name
    new_readings = READINGS_DIR / cleaned_new
    if old_readings.exists() and not new_readings.exists():
        old_readings.rename(new_readings)
    else:
        new_readings.mkdir(parents=True, exist_ok=True)
        
    return cleaned_new

def delete_subject(name):
    import shutil
    classes_path = DB_DIR / name
    readings_path = READINGS_DIR / name
    if classes_path.exists() and classes_path.is_dir():
        shutil.rmtree(classes_path, ignore_errors=True)
    if readings_path.exists() and readings_path.is_dir():
        shutil.rmtree(readings_path, ignore_errors=True)

# Helper para obtener la materia activa de forma persistente
def get_active_subject():
    if "active_subject" not in st.session_state:
        subjects = list_subjects()
        st.session_state.active_subject = subjects[0]
    # Si la materia activa en session_state no está en la lista de materias, reasignar
    subjects = list_subjects()
    if st.session_state.active_subject not in subjects:
        st.session_state.active_subject = subjects[0]
    return st.session_state.active_subject

def get_classes_dir():
    subject = get_active_subject()
    d = DB_DIR / subject
    d.mkdir(parents=True, exist_ok=True)
    return d

def get_readings_dir():
    subject = get_active_subject()
    d = READINGS_DIR / subject
    d.mkdir(parents=True, exist_ok=True)
    return d

# Clases Helpers
def list_saved_classes():
    active_sub = get_active_subject()
    class_names = set()
    act_clean = active_sub.strip().lower()
    
    # 1. Supabase DB
    try:
        from supabase_client import get_supabase_client
        client = get_supabase_client()
        res = client.table("documents").select("title").eq("file_type", "class").execute()
        if res.data:
            for row in res.data:
                t = row.get("title", "")
                if t.startswith("[") and "]" in t:
                    subj = t[1:t.index("]")].strip().lower()
                    name = t[t.index("]") + 1:].strip()
                    if subj == act_clean and name:
                        class_names.add(name)
    except Exception as e:
        print(f"[Supabase list_saved_classes Notice] {e}")

    # 2. Disco local
    classes_dir = get_classes_dir()
    if classes_dir.exists():
        for p in classes_dir.glob("*.json"):
            class_names.add(p.stem)
            
    return sorted(list(class_names))

def load_class_data(name):
    active_sub = get_active_subject()
    doc_title = f"[{active_sub}] {name}"
    
    # 1. Supabase DB
    try:
        from supabase_client import get_supabase_client
        client = get_supabase_client()
        res = client.table("documents").select("*").eq("title", doc_title).execute()
        if res.data and len(res.data) > 0:
            row = res.data[0]
            return {
                "name": name,
                "summary": row.get("summary_markdown", ""),
                "docs_extracted": "",
                "depth": "Detallado",
                "date": str(row.get("created_at", ""))
            }
    except Exception as e:
        print(f"[Supabase load_class_data Notice] {e}")

    # 2. Disco local
    classes_dir = get_classes_dir()
    path = classes_dir / f"{name}.json"
    if path.exists():
        with open(path, 'r', encoding='utf-8') as f:
            return json.load(f)
            
    return None

def save_class_data(name, summary, docs_extracted, depth):
    classes_dir = get_classes_dir()
    path = classes_dir / f"{name}.json"
    data = {
        "name": name,
        "summary": summary,
        "docs_extracted": docs_extracted,
        "depth": depth,
        "date": time.strftime("%Y-%m-%d %H:%M:%S")
    }
    with open(path, 'w', encoding='utf-8') as f:
        json.dump(data, f, ensure_ascii=False, indent=4)

    # Copia en carpeta dedicada de Respaldos Locales (backups_locales)
    try:
        backup_dir = Path(__file__).parent / "backups_locales" / get_active_subject()
        backup_dir.mkdir(parents=True, exist_ok=True)
        with open(backup_dir / f"{name}.json", 'w', encoding='utf-8') as f:
            json.dump(data, f, ensure_ascii=False, indent=4)
    except Exception as e:
        print(f"[Backup Notice] {e}")

    # Sincronización automática con Supabase DB y Vectores RAG
    try:
        from document_service import create_document, update_document_summary
        from rag_service import chunk_text, insert_document_chunks
        active_sub = get_active_subject()
        doc_title = f"[{active_sub}] {name}"
        doc = create_document(title=doc_title, file_type="class")
        doc_id = doc["id"]
        update_document_summary(doc_id, summary, status="completed")
        full_content = f"{doc_title}\n\n{summary}\n\n{docs_extracted}"
        chunks = chunk_text(full_content, chunk_size=800, overlap=100)
        if chunks:
            insert_document_chunks(doc_id, chunks)
    except Exception as e:
        print(f"[Supabase Sync Notice] {e}")

# Lecturas Helpers
def list_saved_readings():
    active_sub = get_active_subject()
    reading_names = set()
    act_clean = active_sub.strip().lower()
    
    # 1. Supabase DB
    try:
        from supabase_client import get_supabase_client
        client = get_supabase_client()
        res = client.table("documents").select("title").eq("file_type", "reading").execute()
        if res.data:
            for row in res.data:
                t = row.get("title", "")
                if t.startswith("[") and "]" in t:
                    subj = t[1:t.index("]")].strip().lower()
                    name = t[t.index("]") + 1:].strip()
                    if subj == act_clean and name:
                        reading_names.add(name)
    except Exception as e:
        print(f"[Supabase list_saved_readings Notice] {e}")

    # 2. Disco local
    readings_dir = get_readings_dir()
    if readings_dir.exists():
        for p in readings_dir.glob("*.json"):
            reading_names.add(p.stem)
            
    return sorted(list(reading_names))

def load_reading_data(name):
    active_sub = get_active_subject()
    doc_title = f"[{active_sub}] {name}"
    
    # 1. Supabase DB
    try:
        from supabase_client import get_supabase_client
        client = get_supabase_client()
        res = client.table("documents").select("*").eq("title", doc_title).execute()
        if res.data and len(res.data) > 0:
            row = res.data[0]
            summary_txt = row.get("summary_markdown", "")
            return {
                "name": name,
                "filename": f"{name}.pdf",
                "pages": [{"page_num": 1, "content": summary_txt}],
                "date": str(row.get("created_at", ""))
            }
    except Exception as e:
        print(f"[Supabase load_reading_data Notice] {e}")

    # 2. Disco local
    readings_dir = get_readings_dir()
    path = readings_dir / f"{name}.json"
    if path.exists():
        with open(path, 'r', encoding='utf-8') as f:
            return json.load(f)
            
    return None

def save_reading_data(name, filename, pages):
    readings_dir = get_readings_dir()
    path = readings_dir / f"{name}.json"
    data = {
        "name": name,
        "filename": filename,
        "pages": pages,
        "date": time.strftime("%Y-%m-%d %H:%M:%S")
    }
    with open(path, 'w', encoding='utf-8') as f:
        json.dump(data, f, ensure_ascii=False, indent=4)

    # Copia en carpeta dedicada de Respaldos Locales (backups_locales)
    try:
        backup_dir = Path(__file__).parent / "backups_locales" / get_active_subject()
        backup_dir.mkdir(parents=True, exist_ok=True)
        with open(backup_dir / f"{name}.json", 'w', encoding='utf-8') as f:
            json.dump(data, f, ensure_ascii=False, indent=4)
    except Exception as e:
        print(f"[Backup Notice] {e}")

    # Sincronización automática con Supabase DB y Vectores RAG
    try:
        from document_service import create_document, update_document_summary
        from rag_service import chunk_text, insert_document_chunks
        active_sub = get_active_subject()
        doc_title = f"[{active_sub}] {name}"
        doc = create_document(title=doc_title, file_type="reading")
        doc_id = doc["id"]
        reading_text = f"Lectura: {name}\n" + "\n".join([p.get('content', '') for p in pages if isinstance(p, dict)])
        update_document_summary(doc_id, reading_text, status="completed")
        chunks = chunk_text(reading_text, chunk_size=800, overlap=100)
        if chunks:
            insert_document_chunks(doc_id, chunks)
    except Exception as e:
        print(f"[Supabase Sync Notice] {e}")

def get_reading_consolidated_text(r_data, max_chars=80000):
    compiled_pages = []
    current_len = 0
    for p in r_data.get('pages', []):
        p_text = f"--- PÁGINA {p['page_num']} ---\n{p['content']}\n\n"
        if current_len + len(p_text) > max_chars:
            compiled_pages.append(f"\n... [Texto de la lectura truncado para optimizar el tamaño de la consulta de IA] ...\n")
            break
        compiled_pages.append(p_text)
        current_len += len(p_text)
    return "".join(compiled_pages)

def get_relevant_reading_pages_context(query, max_pages=15):
    readings = list_saved_readings()
    # Stopwords en español
    stopwords = {"de", "la", "que", "el", "en", "y", "a", "los", "del", "se", "las", "por", "un", "para", "con", "no", "una", "su", "al", "lo", "como", "más", "pero", "sus", "este", "o", "sus", "este", "ese", "si", "ella", "ellos"}
    # Tokenizar query
    query_words = [w.lower() for w in re.findall(r'\w+', query) if w.lower() not in stopwords and len(w) > 2]
    
    if not query_words:
        query_words = [w.lower() for w in re.findall(r'\w+', query) if len(w) > 1]
        
    scored_pages = []
    
    for rd in readings:
        r_data = load_reading_data(rd)
        if not r_data:
            continue
        for p in r_data.get('pages', []):
            content_lower = p['content'].lower()
            score = 0
            for word in query_words:
                count = content_lower.count(word)
                if count > 0:
                    score += count * 5
                    score += 10
            if score > 0:
                scored_pages.append({
                    "reading_name": rd,
                    "page_num": p['page_num'],
                    "content": p['content'],
                    "score": score
                })
                
    scored_pages = sorted(scored_pages, key=lambda x: x['score'], reverse=True)
    top_pages = scored_pages[:max_pages]
    
    if not top_pages:
        fallback_text = []
        for rd in readings[:1]:
            r_data = load_reading_data(rd)
            if r_data:
                for p in r_data.get('pages', [])[:max_pages]:
                    fallback_text.append(f"### LECTURA DE APOYO: {rd} (Pág. {p['page_num']})\n{p['content']}")
        return "\n\n=========================================\n\n".join(fallback_text) if fallback_text else "No hay material de estudio cargado."
        
    formatted_context = []
    for p in top_pages:
        formatted_context.append(f"### LECTURA DE APOYO: {p['reading_name']} (Pág. {p['page_num']}) [Relevancia: {p['score']}]\n{p['content']}")
        
    return "\n\n=========================================\n\n".join(formatted_context)

# --- Helpers de extracción de contenido ---

# Helper para extraer texto y notas de orador de PowerPoint (.pptx)
def extract_text_from_pptx(pptx_path, max_pages=None):
    prs = Presentation(pptx_path)
    slides_data = []
    
    slides_to_read = prs.slides[:max_pages] if max_pages is not None else prs.slides
    
    for i, slide in enumerate(slides_to_read, start=1):
        slide_text_parts = []
        
        # 1. Extraer texto de formas de texto estándar
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text_parts.append(shape.text.strip())
                
            # Extraer texto dentro de tablas si existen en la diapositiva
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_text:
                        slide_text_parts.append(" | ".join(row_text))
                        
        # 2. Extraer notas del orador
        notes_text = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            
        slides_data.append({
            "slide_num": i,
            "content": "\n".join(slide_text_parts),
            "notes": notes_text
        })
        
    return slides_data

# Helper para extraer texto de PDF (.pdf)
def extract_text_from_pdf(pdf_path, max_pages=None):
    reader = PdfReader(pdf_path)
    slides_data = []
    
    pages_to_read = reader.pages[:max_pages] if max_pages is not None else reader.pages
    
    for i, page in enumerate(pages_to_read, start=1):
        text = page.extract_text()
        slides_data.append({
            "slide_num": i,
            "content": text.strip() if text else "",
            "notes": ""  # Los PDFs no tienen notas de orador nativas
        })
        
    return slides_data

# Helper para compilar páginas de una lectura
def extract_reading_pages(temp_path, suffix, max_pages=None):
    if suffix == ".pdf":
        raw_parts = extract_text_from_pdf(temp_path, max_pages=max_pages)
    else:
        raw_parts = extract_text_from_pptx(temp_path, max_pages=max_pages)
        
    pages = []
    for p in raw_parts:
        content = p["content"]
        if "notes" in p and p["notes"]:
            content += f"\n\n[Notas de diapositiva]:\n{p['notes']}"
        pages.append({
            "page_num": p["slide_num"],
            "content": content
        })
    return pages

# Helper para extraer audio de video usando moviepy
def extract_audio_from_video(video_path, output_audio_path):
    video = VideoFileClip(video_path)
    if video.audio is None:
        video.close()
        raise ValueError("El video subido no contiene pista de audio.")
    
    # Escribir el archivo de audio (usamos bitrate de 128k para optimizar)
    video.audio.write_audiofile(output_audio_path, bitrate="128k", logger=None)
    video.close()

# Helper para eliminar archivos temporales de forma segura en Windows sin lanzar PermissionError
def safe_remove_temp_file(file_path):
    if not file_path:
        return
    try:
        if os.path.exists(file_path):
            os.unlink(file_path)
    except Exception as e:
        print(f"[Temp Cleanup Notice] {e}")

# Helper para transformar sub-encabezados de nivel 4 (####) en viñetas destacadas (🔹 **Subtema**)
def clean_markdown_h4_headers(text):
    if not text:
        return text
    lines = text.split('\n')
    cleaned_lines = []
    for line in lines:
        stripped_line = line.strip()
        if stripped_line.startswith('#### '):
            content = stripped_line[5:].strip()
            if content.startswith('**') and content.endswith('**'):
                cleaned_lines.append(f"🔹 {content}")
            else:
                cleaned_lines.append(f"🔹 **{content}**")
        elif stripped_line.startswith('##### ') or stripped_line.startswith('###### '):
            content = stripped_line.lstrip('#').strip()
            if content.startswith('**') and content.endswith('**'):
                cleaned_lines.append(f"▪ {content}")
            else:
                cleaned_lines.append(f"▪ **{content}**")
        else:
            cleaned_lines.append(line)
    return '\n'.join(cleaned_lines)

# Helper para limpiar texto no compatible con codificación Latin-1 de FPDF
def clean_line_for_pdf(line):
    # Reemplazar caracteres de diagramas/cajas unicode por sus equivalentes ASCII antes de filtrar
    replacements = {
        '─': '-',
        '│': '|',
        '┌': '+',
        '┐': '+',
        '└': '+',
        '┘': '+',
        '├': '+',
        '┤': '+',
        '┬': '+',
        '┴': '+',
        '┼': '+',
        '▲': '^',
        '▼': 'v',
        '◄': '<',
        '►': '>',
        # También otros caracteres de diagramas de caja comunes (dobles o gruesos)
        '═': '=',
        '║': '|',
        '╔': '+',
        '╗': '+',
        '╚': '+',
        '╝': '+',
        '╠': '+',
        '╣': '+',
        '╦': '+',
        '╩': '+',
        '╬': '+',
        # Flechas de bloque
        '➔': '->',
        '→': '->',
        '←': '<-',
        '↑': '^',
        '↓': 'v',
        # Iconos de viñetas y subencabezados
        '🔹': '  ' + chr(149) + ' ',
        '▪': '  ' + chr(149) + ' ',
        '▪️': '  ' + chr(149) + ' ',
        '📌': '  ' + chr(149) + ' ',
        '💡': '  ' + chr(149) + ' ',
        '⚡': '  ' + chr(149) + ' ',
        '🎉': '  ' + chr(149) + ' ',
    }
    for unicode_char, ascii_char in replacements.items():
        line = line.replace(unicode_char, ascii_char)
        
    cleaned = "".join(c for c in line if ord(c) < 256)
    cleaned = cleaned.replace('####  ', '#### ')
    cleaned = cleaned.replace('###  ', '### ')
    cleaned = cleaned.replace('##  ', '## ')
    cleaned = cleaned.replace('#  ', '# ')
    cleaned = cleaned.replace('*  ', '* ')
    cleaned = cleaned.replace('-  ', '- ')
    return cleaned

# Helper para estimar la cantidad de líneas que tomará un texto en una celda PDF de cierto ancho
def estimate_pdf_lines(pdf, text, width):
    if not text:
        return 1
    # Estimación de caracteres promedio por línea basándonos en el tamaño de fuente
    char_width_est = pdf.font_size * 0.45
    max_chars = max(1, int(width / char_width_est))
    
    lines = 0
    paragraphs = text.split('\n')
    for p in paragraphs:
        words = p.split(' ')
        current_len = 0
        p_lines = 1
        for w in words:
            if current_len + len(w) + 1 > max_chars:
                p_lines += 1
                current_len = len(w)
            else:
                current_len += len(w) + 1
        lines += p_lines
    return lines

def clean_latex_math(text):
    if not text:
        return text
    # 1. Reemplazar \rightarrow por flecha compatible con Latin-1 ( -> )
    text = text.replace(r"\rightarrow", " -> ")
    text = text.replace(r"\leftarrow", " <- ")
    text = text.replace(r"\to", " -> ")
    
    # 2. Remover \text{...} y extraer su contenido
    for _ in range(5):
        text = re.sub(r'\\text\{([^{}]*)\}', r'\1', text)
        
    # 3. Remover delimitadores de ecuación $$ y \[ \] y \( \)
    text = text.replace("$$", "")
    text = text.replace("\\[", "")
    text = text.replace("\\]", "")
    text = text.replace("\\(", "")
    text = text.replace("\\)", "")
    
    return text

# Helper para generar PDF formateado a partir de Markdown
def generate_pdf_from_markdown(markdown_text):
    markdown_text = clean_markdown_h4_headers(clean_latex_math(markdown_text))
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    
    lines = markdown_text.split('\n')
    
    in_code_block = False
    is_table = False
    table_row_count = 0
    last_line_was_empty = False
    
    for line in lines:
        safe_line = clean_line_for_pdf(line)
        safe_line = safe_line.rstrip()
        
        # Ignorar líneas vacías consecutivas para limpiar espacios en blanco innecesarios
        if safe_line.strip() == "":
            if not last_line_was_empty:
                pdf.ln(3)
                last_line_was_empty = True
            is_table = False
            continue
            
        last_line_was_empty = False
        
        # --- A. Control de bloques de código (Diagramas y esquemas de texto) ---
        if safe_line.startswith('```'):
            in_code_block = not in_code_block
            pdf.ln(2)
            is_table = False
            continue
            
        if in_code_block:
            pdf.set_font("Courier", size=8.5)
            # Dibujar un fondo gris muy claro para simular bloque de código
            current_y = pdf.get_y()
            pdf.set_fill_color(248, 250, 252) # #f8fafc
            pdf.rect(10, current_y, 190, 5.5, 'F')
            pdf.set_text_color(15, 23, 42) # #0f172a
            pdf.set_xy(10, current_y)
            pdf.write(5.5, "    " + safe_line)
            pdf.ln(5.5)
            pdf.set_text_color(0, 0, 0)
            is_table = False
            continue
            
        # --- B. Control de tablas (Líneas con caracter tubería "|") ---
        if '|' in safe_line and not (safe_line.startswith('#') or safe_line.startswith('*') or safe_line.startswith('-')):
            raw_cols = safe_line.split('|')
            if safe_line.startswith('|') and safe_line.endswith('|'):
                cols = [c.strip() for c in raw_cols[1:-1]]
            else:
                cols = [c.strip() for c in raw_cols if c.strip()]
                
            # Saltar la fila separadora de Markdown: |---|---|
            if cols and all(re.match(r'^[\s:-|-]+$', c) for c in cols):
                continue
                
            # Saltar filas completamente vacías: |  |  |
            if not cols or all(c == "" for c in cols):
                continue
                
            num_cols = len(cols)
            if num_cols > 0:
                col_width = 180 / num_cols
                
                # Control de salto de página inteligente
                start_y = pdf.get_y()
                if start_y > 240:
                    pdf.add_page()
                    start_y = pdf.get_y()
                
                # Paso 1: Medir la altura de cada columna en la fila
                max_cell_height = 6 # Altura mínima de fila
                pdf.set_font("Helvetica", '', 9.5) # Establecer font temporal para medición
                for col in cols:
                    col_cleaned = col.replace('**', '').strip()
                    num_lines = estimate_pdf_lines(pdf, col_cleaned, col_width)
                    cell_height = num_lines * 5.5 + 2 # Altura de texto más padding
                    if cell_height > max_cell_height:
                        max_cell_height = cell_height
                
                # Paso 2: Dibujar cabecera o datos con bordes alineados y colores
                if not is_table:
                    is_table = True
                    table_row_count = 0
                else:
                    table_row_count += 1
                
                for i, col in enumerate(cols):
                    col_x = 10 + i * col_width
                    col_cleaned = col.replace('**', '').strip()
                    
                    if table_row_count == 0:
                        # Cabecera: fondo morado-azul, texto blanco en negrita
                        pdf.set_fill_color(99, 102, 241) # #6366f1
                        pdf.set_text_color(255, 255, 255)
                        pdf.set_font("Helvetica", 'B', 9)
                        pdf.rect(col_x, start_y, col_width, max_cell_height, 'FD')
                    else:
                        # Datos: alternar fondo blanco y gris claro (Zebra striping)
                        if table_row_count % 2 == 1:
                            pdf.set_fill_color(248, 250, 252) # #f8fafc
                            pdf.rect(col_x, start_y, col_width, max_cell_height, 'FD')
                        else:
                            pdf.set_fill_color(255, 255, 255)
                            pdf.rect(col_x, start_y, col_width, max_cell_height, 'FD')
                        
                        pdf.set_text_color(30, 41, 59) # Gris oscuro para texto
                        pdf.set_font("Helvetica", '', 9)
                    
                    # Dibujar texto
                    pdf.set_xy(col_x, start_y + 1)
                    pdf.set_draw_color(203, 213, 225) # #cbd5e1 Borde gris claro
                    pdf.multi_cell(col_width, 5, col_cleaned, border=0, align='L')
                    
                # Restaurar colores por defecto
                pdf.set_draw_color(0, 0, 0)
                pdf.set_text_color(0, 0, 0)
                
                pdf.set_y(start_y + max_cell_height)
                continue
        else:
            is_table = False
            
        # --- C. Elementos estándar de Markdown ---
        
        # 1. H1 Header
        if safe_line.startswith('# '):
            header_text = safe_line[2:].replace('**', '').strip()
            pdf.ln(4)
            pdf.set_font("Helvetica", 'B', 16)
            pdf.multi_cell(0, 8, header_text)
            pdf.ln(4)
            
        # 2. H2 Header
        elif safe_line.startswith('## '):
            header_text = safe_line[3:].replace('**', '').strip()
            pdf.ln(3)
            pdf.set_font("Helvetica", 'B', 14)
            pdf.multi_cell(0, 7, header_text)
            pdf.ln(3)
            
        # 3. H3 Header
        elif safe_line.startswith('### '):
            header_text = safe_line[4:].replace('**', '').strip()
            pdf.ln(2)
            pdf.set_font("Helvetica", 'B', 12)
            pdf.multi_cell(0, 6, header_text)
            pdf.ln(2)
            
        # 4. Línea Horizontal (markdown ---)
        elif safe_line == "---" or safe_line == "___":
            pdf.ln(2)
            pdf.line(pdf.get_x(), pdf.get_y(), pdf.get_x() + 180, pdf.get_y())
            pdf.ln(4)
            
        # 5. Listas con viñetas (* o -)
        elif safe_line.startswith('* ') or safe_line.startswith('- '):
            pdf.set_font("Helvetica", size=10)
            pdf.write(6, chr(149) + " ")
            clean_list_line = re.sub(r'(?<!\*)\*(?!\*)', '', safe_line[2:])
            parts = clean_list_line.split('**')
            is_bold = False
            for part in parts:
                if is_bold:
                    pdf.set_font("Helvetica", "B", 10)
                else:
                    pdf.set_font("Helvetica", "", 10)
                pdf.write(6, part)
                is_bold = not is_bold
            pdf.ln(8)
            
        # 6. Listas numeradas (1. , 2. etc.)
        elif re.match(r'^(\d+\.\s)(.*)', safe_line):
            match = re.match(r'^(\d+\.\s)(.*)', safe_line)
            prefix = match.group(1)
            rest = match.group(2)
            pdf.set_font("Helvetica", size=10)
            pdf.write(6, prefix)
            clean_num_line = re.sub(r'(?<!\*)\*(?!\*)', '', rest)
            parts = clean_num_line.split('**')
            is_bold = False
            for part in parts:
                if is_bold:
                    pdf.set_font("Helvetica", "B", 10)
                else:
                    pdf.set_font("Helvetica", "", 10)
                pdf.write(6, part)
                is_bold = not is_bold
            pdf.ln(8)

        # 7. Citas y Notas Destacadas (Blockquotes: líneas que inician con >)
        elif safe_line.startswith('>'):
            quote_text = safe_line.lstrip('>').strip()
            # Eliminar asteriscos individuales de cursiva (* "..." *) preservando los de negrita (**)
            quote_text = re.sub(r'(?<!\*)\*(?!\*)', '', quote_text).strip()
            
            start_y = pdf.get_y()
            if start_y > 260:
                pdf.add_page()
                start_y = pdf.get_y()
                
            pdf.set_font("Helvetica", "I", 9.5)
            pdf.set_text_color(51, 65, 85) # #334155 Slate oscuro
            
            # Escribir el texto sangrado a la derecha
            pdf.set_xy(16, start_y)
            parts = quote_text.split('**')
            is_bold = False
            for part in parts:
                if is_bold:
                    pdf.set_font("Helvetica", "BI", 9.5)
                else:
                    pdf.set_font("Helvetica", "I", 9.5)
                pdf.write(5.5, part)
                is_bold = not is_bold
                
            end_y = pdf.get_y()
            block_height = max(6, end_y - start_y)
            
            # Dibujar barra vertical morada en el margen izquierdo
            pdf.set_draw_color(99, 102, 241) # #6366f1
            pdf.set_line_width(0.8)
            pdf.line(12, start_y, 12, start_y + block_height)
            pdf.set_line_width(0.2) # Restaurar ancho de línea
            
            # Restaurar colores por defecto
            pdf.set_draw_color(0, 0, 0)
            pdf.set_text_color(0, 0, 0)
            pdf.ln(7)
            
        # 8. Párrafos normales con soporte para texto en negrita
        else:
            safe_paragraph = re.sub(r'(?<!\*)\*(?!\*)', '', safe_line)
            pdf.set_font("Helvetica", size=10)
            parts = safe_paragraph.split('**')
            is_bold = False
            for part in parts:
                if is_bold:
                    pdf.set_font("Helvetica", "B", 10)
                else:
                    pdf.set_font("Helvetica", "", 10)
                pdf.write(6, part)
                is_bold = not is_bold
            pdf.ln(8)
            
    return bytes(pdf.output())

# Helper para generar PDF con diseño especial de "Fácil Lectura" (Configurable)
def generate_easy_read_pdf(text, font_family="Sans-Serif", font_size=11.5, line_spacing=1.5, alignment="Justificado"):
    text = clean_markdown_h4_headers(clean_latex_math(text))
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=20)
    pdf.add_page()
    
    # Mapear familia de fuentes
    pdf_font = "Helvetica"
    if "Serif" in font_family:
        pdf_font = "Times"
    elif "Monospace" in font_family:
        pdf_font = "Courier"
        
    lines = text.split('\n')
    
    # Calcular altura de celda
    h_cell = font_size * 0.65 * line_spacing
    pdf_align = "J" if alignment == "Justificado" else "L"
    
    in_code_block = False
    is_table = False
    table_row_count = 0
    last_line_was_empty = False
    
    for line in lines:
        safe_line = clean_line_for_pdf(line)
        safe_line = safe_line.rstrip()
        
        # Ignorar líneas vacías consecutivas para limpiar espacios en blanco innecesarios
        if safe_line.strip() == "":
            if not last_line_was_empty:
                pdf.ln(h_cell * 0.5)
                last_line_was_empty = True
            is_table = False
            continue
            
        last_line_was_empty = False
        
        # --- A. Control de bloques de código (Diagramas y esquemas de texto) ---
        if safe_line.startswith('```'):
            in_code_block = not in_code_block
            pdf.ln(2)
            is_table = False
            continue
            
        if in_code_block:
            pdf.set_font("Courier", size=font_size * 0.8)
            # Dibujar un fondo gris muy claro para simular bloque de código
            current_y = pdf.get_y()
            pdf.set_fill_color(248, 250, 252) # #f8fafc
            pdf.rect(10, current_y, 190, h_cell * 0.8, 'F')
            pdf.set_text_color(15, 23, 42) # #0f172a
            pdf.set_xy(10, current_y)
            pdf.write(h_cell * 0.8, "    " + safe_line)
            pdf.ln(h_cell * 0.8)
            pdf.set_text_color(0, 0, 0)
            is_table = False
            continue
            
        # --- B. Control de tablas (Líneas con caracter tubería "|") ---
        if '|' in safe_line and not (safe_line.startswith('#') or safe_line.startswith('*') or safe_line.startswith('-')):
            raw_cols = safe_line.split('|')
            if safe_line.startswith('|') and safe_line.endswith('|'):
                cols = [c.strip() for c in raw_cols[1:-1]]
            else:
                cols = [c.strip() for c in raw_cols if c.strip()]
                
            # Saltar la fila separadora de Markdown: |---|---|
            if cols and all(re.match(r'^[\s:-|-]+$', c) for c in cols):
                continue
                
            # Saltar filas completamente vacías: |  |  |
            if not cols or all(c == "" for c in cols):
                continue
                
            num_cols = len(cols)
            if num_cols > 0:
                col_width = 180 / num_cols
                
                # Control de salto de página inteligente
                start_y = pdf.get_y()
                if start_y > 230:
                    pdf.add_page()
                    start_y = pdf.get_y()
                
                # Paso 1: Medir la altura de cada columna en la fila
                max_cell_height = font_size * 0.6 * line_spacing + 2 # Altura mínima de fila
                pdf.set_font(pdf_font, '', font_size * 0.8) # Establecer font temporal para medición
                for col in cols:
                    col_cleaned = col.replace('**', '').strip()
                    num_lines = estimate_pdf_lines(pdf, col_cleaned, col_width)
                    cell_height = num_lines * (font_size * 0.8 * 0.7 * line_spacing) + 2
                    if cell_height > max_cell_height:
                        max_cell_height = cell_height
                
                # Paso 2: Dibujar cabecera o datos con bordes alineados y colores
                if not is_table:
                    is_table = True
                    table_row_count = 0
                else:
                    table_row_count += 1
                
                for i, col in enumerate(cols):
                    col_x = 10 + i * col_width
                    col_cleaned = col.replace('**', '').strip()
                    
                    if table_row_count == 0:
                        # Cabecera: fondo morado-azul, texto blanco en negrita
                        pdf.set_fill_color(99, 102, 241) # #6366f1
                        pdf.set_text_color(255, 255, 255)
                        pdf.set_font(pdf_font, 'B', font_size * 0.8)
                        pdf.rect(col_x, start_y, col_width, max_cell_height, 'FD')
                    else:
                        # Datos: alternar fondo blanco y gris claro
                        if table_row_count % 2 == 1:
                            pdf.set_fill_color(248, 250, 252) # #f8fafc
                            pdf.rect(col_x, start_y, col_width, max_cell_height, 'FD')
                        else:
                            pdf.set_fill_color(255, 255, 255)
                            pdf.rect(col_x, start_y, col_width, max_cell_height, 'FD')
                        
                        pdf.set_text_color(30, 41, 59)
                        pdf.set_font(pdf_font, '', font_size * 0.8)
                    
                    # Dibujar texto
                    pdf.set_xy(col_x, start_y + 1)
                    pdf.set_draw_color(203, 213, 225)
                    pdf.multi_cell(col_width, (font_size * 0.8 * 0.7 * line_spacing), col_cleaned, border=0, align='L')
                    
                pdf.set_draw_color(0, 0, 0)
                pdf.set_text_color(0, 0, 0)
                
                pdf.set_y(start_y + max_cell_height)
                continue
        else:
            is_table = False
            
        # --- C. Elementos estándar de Markdown ---
        
        # 1. H1 Header
        if safe_line.startswith('# '):
            header_text = safe_line[2:].replace('**', '').strip()
            pdf.ln(font_size * 0.5)
            pdf.set_font(pdf_font, 'B', font_size * 1.5)
            pdf.multi_cell(0, font_size * 1.5 * 0.7, header_text, align=pdf_align)
            pdf.ln(font_size * 0.5)
        elif safe_line.startswith('## '):
            header_text = safe_line[3:].replace('**', '').strip()
            pdf.ln(font_size * 0.4)
            pdf.set_font(pdf_font, 'B', font_size * 1.3)
            pdf.multi_cell(0, font_size * 1.3 * 0.7, header_text, align=pdf_align)
            pdf.ln(font_size * 0.4)
        elif safe_line.startswith('### '):
            header_text = safe_line[4:].replace('**', '').strip()
            pdf.ln(font_size * 0.3)
            pdf.set_font(pdf_font, 'B', font_size * 1.1)
            pdf.multi_cell(0, font_size * 1.1 * 0.7, header_text, align=pdf_align)
            pdf.ln(font_size * 0.3)
        elif safe_line.startswith('* ') or safe_line.startswith('- '):
            pdf.set_font(pdf_font, size=font_size)
            clean_bullet = re.sub(r'(?<!\*)\*(?!\*)', '', safe_line[2:]).replace('**', '').strip()
            bullet_text = f"  {chr(149)}  " + clean_bullet
            pdf.multi_cell(0, h_cell, bullet_text, align=pdf_align)
            pdf.ln(2)
        elif safe_line.startswith('>'):
            quote_text = safe_line.lstrip('>').strip()
            quote_text = re.sub(r'(?<!\*)\*(?!\*)', '', quote_text).strip()
            
            start_y = pdf.get_y()
            if start_y > 250:
                pdf.add_page()
                start_y = pdf.get_y()
                
            pdf.set_font(pdf_font, "I", font_size * 0.95)
            pdf.set_text_color(51, 65, 85)
            
            pdf.set_xy(16, start_y)
            parts = quote_text.split('**')
            is_bold = False
            for part in parts:
                if is_bold:
                    pdf.set_font(pdf_font, "BI", font_size * 0.95)
                else:
                    pdf.set_font(pdf_font, "I", font_size * 0.95)
                pdf.write(h_cell, part)
                is_bold = not is_bold
                
            end_y = pdf.get_y()
            block_height = max(h_cell, end_y - start_y)
            
            pdf.set_draw_color(99, 102, 241)
            pdf.set_line_width(0.8)
            pdf.line(12, start_y, 12, start_y + block_height)
            pdf.set_line_width(0.2)
            
            pdf.set_draw_color(0, 0, 0)
            pdf.set_text_color(0, 0, 0)
            pdf.ln(h_cell * 0.5)
        else:
            safe_paragraph = re.sub(r'(?<!\*)\*(?!\*)', '', safe_line).replace('**', '')
            pdf.set_font(pdf_font, size=font_size)
            pdf.multi_cell(0, h_cell, safe_paragraph, align=pdf_align)
            pdf.ln(3)
            
    return bytes(pdf.output())

# Helper para parsear spoilers/respuestas ocultas de cuestionarios y mostrarlos como expanders de Streamlit
def render_study_material_with_expanders(text):
    pattern = r':::spoiler\s*(.*?)\n(.*?)\n:::'
    matches = re.split(pattern, text, flags=re.DOTALL)
    
    i = 0
    while i < len(matches):
        st.markdown(matches[i])
        if i + 2 < len(matches):
            title = matches[i+1].strip()
            content = matches[i+2].strip()
            with st.expander(title):
                st.markdown(content)
            i += 3
        else:
            break

# Helper para reubicar respuestas de cuestionarios al final del documento en descargas
def adjust_quiz_answers_for_download(text):
    spoilers = []
    counter = [1]
    
    def replace_spoiler(match):
        title = match.group(1).strip()
        content = match.group(2).strip()
        spoilers.append((counter[0], title, content))
        replacement = f"\n*(Ver respuesta {counter[0]} al final del documento)*\n"
        counter[0] += 1
        return replacement
        
    pattern = r':::spoiler\s*(.*?)\n(.*?)\n:::'
    modified_text = re.sub(pattern, replace_spoiler, text, flags=re.DOTALL)
    
    if spoilers:
        modified_text += "\n\n---\n\n## 🔑 Hoja de Respuestas\n\n"
        for num, title, content in spoilers:
            modified_text += f"### Respuesta Pregunta {num}\n"
            modified_text += f"{content}\n\n"
            
    return modified_text

# Helper para manejar y dar recomendaciones amigables sobre errores de la API de Gemini (Rate Limits, Quota)
def handle_api_error(e):
    err_str = str(e)
    if "RESOURCE_EXHAUSTED" in err_str or "429" in err_str or "quota" in err_str.lower():
        return (
            "⚠️ **Límite de Cuota Excedido (429):** Has superado el límite de tokens por minuto del nivel gratuito de la API de Gemini (250,000 TPM).\n\n"
            "**Cómo solucionarlo:**\n"
            "1. **Espera 1 o 2 minutos** para que tu cuota de tokens por minuto se reinicie y vuelve a intentarlo.\n"
            "2. **Evita enviar textos gigantescos en una sola consulta** (por ejemplo, libros de más de 30-40 páginas completas).\n"
            "3. **Usa herramientas locales**: En la extracción de capítulos, prefiere usar **'Por Rango de Páginas (Local)'**, el cual es instantáneo, 100% gratuito y no consume nada de cuota de API de Gemini.\n"
            "4. **Aumenta tu cuota**: Puedes pasar tu API Key a facturación mensual (Pay-As-You-Go) en tu consola de Google AI Studio. Esto incrementará tus límites drásticamente a millones de tokens por minuto, y el consumo real es sumamente barato (fracciones de centavo de dólar)."
        )
    if "UNAVAILABLE" in err_str or "503" in err_str:
        return (
            "⚠️ **Servicio No Disponible (503):** El modelo de Gemini seleccionado está experimentando una alta demanda temporal en los servidores de Google.\n\n"
            "**Cómo solucionarlo:**\n"
            "1. **Cambia el modelo en la barra lateral**: En la sección de configuración de la izquierda (debajo del campo de API Key), selecciona un modelo alternativo como **'gemini-2.0-flash'** o **'gemini-3.1-flash-lite'** y vuelve a intentarlo. Esto redirigirá tu petición a otros servidores activos de inmediato.\n"
            "2. **Espera unos segundos** y vuelve a intentar generar tu resumen. La saturación de los servidores suele ser momentánea."
        )
    return f"❌ Error de la API de Gemini: {err_str}"

# Helper para actualizar el overlay de carga en un placeholder
def update_full_screen_loader(placeholder, message):
    placeholder.markdown(f"""
    <div style="
        position: fixed;
        top: 0;
        left: 0;
        width: 100vw;
        height: 100vh;
        background-color: rgba(15, 17, 26, 0.92);
        z-index: 999999;
        display: flex;
        flex-direction: column;
        align-items: center;
        justify-content: center;
        color: #f8fafc;
    ">
        <div style="
            border: 8px solid #334155;
            border-top: 8px solid #8b5cf6;
            border-radius: 50%;
            width: 80px;
            height: 80px;
            animation: spin 1.5s linear infinite;
            margin-bottom: 25px;
        "></div>
        <style>
            @keyframes spin {{
                0% {{ transform: rotate(0deg); }}
                100% {{ transform: rotate(360deg); }}
            }}
        </style>
        <h3 style="font-family: 'Outfit', sans-serif; font-weight: 600; color: #8b5cf6; margin: 0 0 10px 0; font-size: 1.8rem;">ClaseSinc AI</h3>
        <p style="font-family: 'Inter', sans-serif; color: #e2e8f0; font-size: 1.2rem; margin: 0 0 5px 0; text-align: center; font-weight: 500;">
            {message}
        </p>
        <p style="font-family: 'Inter', sans-serif; color: #94a3b8; font-size: 0.95rem; margin: 0; text-align: center;">
            Por favor, espera. No recargues la página ni cambies de pestaña.
        </p>
    </div>
    """, unsafe_allow_html=True)

# Helper para mostrar un overlay de carga a pantalla completa bloqueando cualquier interacción del usuario
def show_full_screen_loader(message="Procesando..."):
    placeholder = st.empty()
    update_full_screen_loader(placeholder, message)
    return placeholder

# Helper para buscar y extraer únicamente las páginas relevantes a una consulta (para evitar rebasar límites de tokens)
def get_relevant_pages_for_extraction(reading_data, query, window_size=35):
    if len(reading_data["pages"]) <= window_size:
        return "\n\n".join([f"--- PÁGINA {p['page_num']} ---\n{p['content']}" for p in reading_data["pages"]])
        
    query_clean = query.lower().strip()
    
    # Extraer palabras únicas de la consulta mayores a 4 caracteres para calificar relevancia
    words = [w for w in re.findall(r'\w+', query_clean) if len(w) > 4]
    if not words:
        words = [w for w in re.findall(r'\w+', query_clean) if len(w) > 2]
        
    page_scores = []
    
    for p in reading_data["pages"]:
        content_lower = p["content"].lower()
        
        # Penalizar fuertemente si es un índice general de la obra completa
        if "índice" in content_lower and "tabla de contenido" in content_lower:
            score = 0
        else:
            # Calcular la frecuencia agregada de palabras clave en la página
            score = sum(content_lower.count(word) for word in words)
            
            # Penalizar fuertemente si la página tiene las palabras "prólogo" o "introducción de la obra"
            if "prólogo" in content_lower or "prologo" in content_lower or "estructura de la obra" in content_lower or "contenido del cd" in content_lower:
                score = score * 0.05
                
        page_scores.append((p["page_num"], score))
        
    # Obtener la página con el puntaje más alto
    best_page, max_score = max(page_scores, key=lambda x: x[1])
    
    # Si la frecuencia de palabras clave es baja, usar el fallback del número de capítulo
    if max_score < 3:
        match = re.search(r'(cap[íi]tulo|cap\.?)\s*(\d+)', query_clean)
        if match:
            cap_num = match.group(2)
            search_terms = [f"capítulo {cap_num}", f"capitulo {cap_num}", f"cap. {cap_num}"]
            for page_num, score in page_scores:
                content_lower = next(p["content"].lower() for p in reading_data["pages"] if p["page_num"] == page_num)
                if any(term in content_lower for term in search_terms):
                    # Filtrar que no sea prólogo tampoco en el fallback
                    if "prólogo" not in content_lower and "prologo" not in content_lower:
                        best_page = page_num
                        break
                        
    # Establecer la ventana de páginas centrada en la página de mayor relevancia
    start_idx = max(1, best_page - 2)
    end_idx = min(len(reading_data["pages"]), start_idx + window_size - 1)
    
    relevant_pages = []
    for p in reading_data["pages"]:
        if start_idx <= p["page_num"] <= end_idx:
            relevant_pages.append(f"--- PÁGINA {p['page_num']} ---\n{p['content']}")
            
    return "\n\n".join(relevant_pages)


def split_audio_file_if_needed(audio_path, max_size_mb=24.0):
    file_size_mb = os.path.getsize(audio_path) / (1024 * 1024)
    if file_size_mb <= max_size_mb:
        return [audio_path]
    
    try:
        from moviepy import AudioFileClip
    except ImportError:
        from moviepy.editor import AudioFileClip
        
    clip = AudioFileClip(audio_path)
    duration = clip.duration
    import math
    num_parts = math.ceil(file_size_mb / max_size_mb)
    part_duration = duration / num_parts
    
    parts_paths = []
    base_dir = os.path.dirname(audio_path)
    base_name = os.path.splitext(os.path.basename(audio_path))[0]
    
    for i in range(num_parts):
        start_time = i * part_duration
        end_time = min((i + 1) * part_duration, duration)
        
        if hasattr(clip, "subclipped"):
            part_clip = clip.subclipped(start_time, end_time)
        else:
            part_clip = clip.subclip(start_time, end_time)
            
        part_path = os.path.join(base_dir, f"{base_name}_part_{i+1}.mp3")
        part_clip.write_audiofile(part_path, bitrate="64k", logger=None)
        parts_paths.append(part_path)
        part_clip.close()
        
    clip.close()
    return parts_paths

# --- Configuración del Sidebar ---
st.sidebar.markdown("<h2 style='text-align: center; color: #8b5cf6;'>⚙️ Configuración</h2>", unsafe_allow_html=True)

# Cargar API Keys por defecto desde .env
api_key = os.getenv("GEMINI_API_KEY", "")
openai_api_key = os.getenv("OPENAI_API_KEY", "")

# Selector de modelo de IA (Google y OpenAI)
model_options = [
    "gemini-3.5-flash",
    "gemini-2.0-flash",
    "gemini-3.1-flash-lite",
    "gemini-flash-latest",
    "gpt-4o",
    "gpt-4o-mini"
]
selected_model = st.sidebar.selectbox(
    "Modelo de IA",
    options=model_options,
    index=0,
    help="Selecciona el modelo de Inteligencia Artificial a utilizar. Se sugiere 'gemini-3.5-flash' o 'gpt-4o'."
)

# Entrada de API Key condicional
if "gpt-" in selected_model:
    openai_api_key = st.sidebar.text_input(
        "OpenAI API Key",
        value=openai_api_key,
        type="password",
        help="Introduce tu clave de API de OpenAI (comienza con sk-...)."
    )
    if not openai_api_key:
        st.sidebar.warning("🔑 Por favor ingresa una API Key de OpenAI válida para continuar.")
else:
    api_key = st.sidebar.text_input(
        "Gemini API Key",
        value=api_key,
        type="password",
        help="Introduce tu clave de API de Gemini."
    )
    if not api_key:
        st.sidebar.warning("🔑 Por favor ingresa una API Key de Gemini válida para continuar.")

# Expander de Estado de Modelos en Tiempo Real
with st.sidebar.expander("🚦 Estado de los Modelos", expanded=False):
    st.markdown("Verifica la disponibilidad en tiempo real de los servidores:")
    
    # Mostrar estado de cada modelo
    for m in model_options:
        status_key = f"status_cache_{m}"
        status_val = st.session_state.get(status_key, "⚪ Sin verificar")
        st.markdown(f"**{m}**:\n`{status_val}`")
        
    if st.button("🔍 Probar Conexión", key="btn_test_model_status", use_container_width=True):
        with st.spinner("Probando conexión con los servidores..."):
            for m in model_options:
                status_key = f"status_cache_{m}"
                try:
                    start_time = time.time()
                    if "gpt-" in m:
                        if not openai_api_key:
                            st.session_state[status_key] = "🔴 Falta OpenAI Key"
                            continue
                        test_openai_client = OpenAI(api_key=openai_api_key)
                        test_openai_client.chat.completions.create(
                            model=m,
                            messages=[{"role": "user", "content": "ping"}],
                            max_tokens=1
                        )
                    else:
                        if not api_key:
                            st.session_state[status_key] = "🔴 Falta Gemini Key"
                            continue
                        test_client = genai.Client(api_key=api_key)
                        test_client.models.generate_content(
                            model=m,
                            contents=["ping"],
                            config=dict(max_output_tokens=1)
                        )
                    latency = int((time.time() - start_time) * 1000)
                    st.session_state[status_key] = f"🟢 Operacional ({latency}ms)"
                except Exception as e:
                    err_str = str(e)
                    if "503" in err_str or "UNAVAILABLE" in err_str:
                        st.session_state[status_key] = "🟡 Sobrecargado (503)"
                    elif "404" in err_str or "NOT_FOUND" in err_str:
                        st.session_state[status_key] = "🔴 No Disponible (404)"
                    elif "RESOURCE_EXHAUSTED" in err_str or "429" in err_str:
                        st.session_state[status_key] = "🟡 Cuota Agotada (429)"
                    else:
                        st.session_state[status_key] = "🔴 Error de Conexión"
        st.rerun()

st.sidebar.markdown("---")
st.sidebar.markdown("<h3 style='color: #8b5cf6;'>📚 Materia de Estudio</h3>", unsafe_allow_html=True)

# Obtener lista de materias
subjects_list = list_subjects()
active_subject = get_active_subject()

# Asegurar que el índice seleccionado por defecto coincida con la materia activa
try:
    subject_index = subjects_list.index(active_subject)
except ValueError:
    subject_index = 0

selected_subject = st.sidebar.selectbox(
    "Materia Activa",
    options=subjects_list,
    index=subject_index,
    help="Todo el contenido (clases, lecturas y zona de estudio) se filtrará por esta materia."
)

if selected_subject != active_subject:
    st.session_state.active_subject = selected_subject
    st.rerun()

# Expander para administrar materias
with st.sidebar.expander("⚙️ Administrar Materias"):
    # 1. Crear Materia
    st.markdown("##### ➕ Crear Nueva Materia")
    new_sub_name = st.text_input("Nombre de la materia", value="", key="create_subject_input")
    if st.button("➕ Crear Materia", key="btn_create_subject"):
        if not new_sub_name.strip():
            st.warning("⚠️ Ingresa un nombre.")
        else:
            created = create_subject(new_sub_name)
            if created:
                st.session_state.active_subject = created
                st.success(f"🎉 Materia '{created}' creada con éxito.")
                time.sleep(1.5)
                st.rerun()
            else:
                st.error("❌ Nombre de materia no válido.")

    # 2. Renombrar Materia
    st.markdown("##### ✏️ Renombrar Materia Activa")
    rename_sub_name = st.text_input("Nuevo nombre", value=selected_subject, key="rename_subject_input")
    if st.button("✏️ Renombrar Materia", key="btn_rename_subject"):
        if not rename_sub_name.strip():
            st.warning("⚠️ El nombre no puede estar vacío.")
        elif rename_sub_name.strip() == selected_subject:
            st.info("💡 El nombre es idéntico.")
        else:
            renamed = rename_subject(selected_subject, rename_sub_name)
            if renamed:
                st.session_state.active_subject = renamed
                st.success(f"✏️ Materia renombrada a '{renamed}'.")
                time.sleep(1.5)
                st.rerun()
            else:
                st.error("❌ Nombre de materia no válido o ya existe.")

    # 3. Eliminar Materia
    st.markdown("##### 🗑️ Eliminar Materia Activa")
    st.write(f"Se eliminará la materia `{selected_subject}` con **todo** su contenido (clases y lecturas asociadas).")
    confirm_del_sub = st.checkbox("Confirmar eliminación irreversible", key="confirm_del_subject_cb")
    if st.button("🗑️ Eliminar Materia", key="btn_del_subject", disabled=not confirm_del_sub):
        # Evitar eliminar si es la única materia
        if len(subjects_list) <= 1 and selected_subject == "General":
            st.error("❌ No puedes eliminar la materia 'General' si es la única disponible.")
        else:
            delete_subject(selected_subject)
            st.success(f"🗑️ Materia '{selected_subject}' y su contenido eliminados con éxito.")
            # Cambiar a otra materia restante
            remaining = [s for s in list_subjects() if s != selected_subject]
            st.session_state.active_subject = remaining[0] if remaining else "General"
            time.sleep(1.5)
            st.rerun()

st.sidebar.markdown("---")
st.sidebar.markdown("<h3 style='color: #8b5cf6;'>📝 Nivel de Detalle</h3>", unsafe_allow_html=True)
summary_depth = st.sidebar.selectbox(
    "Profundidad del Resumen",
    options=[
        "Completo y Detallado (Estudio exhaustivo)",
        "Estándar (Equilibrado)",
        "Acotado (Conciso y directo)"
    ],
    index=0,
    help="Determina el nivel de detalle del resumen generado al procesar una nueva clase."
)

st.sidebar.markdown("---")
st.sidebar.markdown("<h3 style='color: #8b5cf6;'>💾 Descarga</h3>", unsafe_allow_html=True)
output_filename = st.sidebar.text_input(
    "Nombre del archivo",
    value="resumen_clase_sincronizado",
    help="Indica el nombre de archivo con el que se descargará el resumen (sin extensión)."
)

st.sidebar.markdown("---")
st.sidebar.markdown("""
### 💡 Instrucciones de Uso:
1. **Configura tu API Key** (enlazada automáticamente al `.env` por defecto).
2. **Elige la pestaña deseada** en el menú principal para procesar, consultar el historial, estudiar o gestionar lecturas.
3. **Guardado Automático**: Al procesar una clase en la pestaña 1, esta quedará almacenada localmente.
""")

# --- Contenido Principal ---
st.markdown("<h1 class='title-gradient'>ClaseSinc AI 🎓</h1>", unsafe_allow_html=True)
st.markdown("<p class='subtitle-text'>Sincronizador inteligente de grabaciones de voz, materiales de apoyo y lecturas complementarias</p>", unsafe_allow_html=True)

# Crear pestañas superiores de navegación
tab_process, tab_history, tab_study, tab_readings, tab_chat = st.tabs([
    "🎙️ Procesar Nueva Clase", 
    "📚 Historial de Clases", 
    "🎓 Zona de Estudio e Integración",
    "📖 Lecturas de Apoyo",
    "💬 Tutor IA (Chat)"
])

# Variables de Session State para persistir salidas durante la ejecución
if "summary_output" not in st.session_state:
    st.session_state.summary_output = None
if "docs_extracted" not in st.session_state:
    st.session_state.docs_extracted = None
if "study_material" not in st.session_state:
    st.session_state.study_material = None
if "study_title" not in st.session_state:
    st.session_state.study_title = ""
if "consolidated_material" not in st.session_state:
    st.session_state.consolidated_material = None
if "consolidated_title" not in st.session_state:
    st.session_state.consolidated_title = ""
if "reading_summary" not in st.session_state:
    st.session_state.reading_summary = None
if "reading_easy_text" not in st.session_state:
    st.session_state.reading_easy_text = None
if "extracted_chapter" not in st.session_state:
    st.session_state.extracted_chapter = None
if "reading_last_action" not in st.session_state:
    st.session_state.reading_last_action = ""

# =====================================================================
# PESTAÑA 1: PROCESAR NUEVA CLASE
# =====================================================================
with tab_process:
    upload_mode = st.radio(
        "Elige la forma de ingresar la clase:",
        options=[
            "🎙️ Procesar desde grabación (Audio/Video + PPTX/PDF)",
            "📥 Importar resumen existente (Archivo .md o .pdf)"
        ],
        horizontal=True
    )
    
    if "Importar resumen" in upload_mode:
        st.markdown("#### 📥 Importar Resumen desde Archivo (.md o .pdf)")
        
        imported_file = st.file_uploader(
            "Selecciona el archivo de resumen (.md, .pdf)",
            type=["md", "pdf"],
            key="imported_file_uploader",
            help="Sube un archivo de resumen existente en formato Markdown (.md) o PDF (.pdf)."
        )
        
        # Inicializar el estado de sesión para el nombre importado si no existe
        if "import_class_name_val" not in st.session_state:
            st.session_state.import_class_name_val = ""
            st.session_state.last_imported_file_name = ""
            
        # Si se sube un nuevo archivo, actualizar la variable de estado
        if imported_file and st.session_state.last_imported_file_name != imported_file.name:
            st.session_state.import_class_name_val = Path(imported_file.name).stem
            st.session_state.last_imported_file_name = imported_file.name
            
        class_name_import = st.text_input(
            "Nombre de la Clase (para el historial)", 
            value=st.session_state.import_class_name_val, 
            placeholder="Ej: Clase 1 - Introduccion al Desarrollo Cognitivo",
            key="import_class_name_input",
            help="Escribe el nombre con el que se guardará esta clase en tu base de datos local."
        )
        
        # Sincronizar el estado con lo ingresado manualmente
        st.session_state.import_class_name_val = class_name_import
        
        if st.button("📥 Importar y Guardar en Historial", key="import_btn"):
            if not imported_file:
                st.warning("⚠️ Por favor selecciona un archivo para importar.")
            elif not class_name_import.strip():
                st.warning("⚠️ Debes ingresar un nombre para la clase.")
            else:
                summary_content = ""
                suffix = Path(imported_file.name).suffix.lower()
                
                with st.spinner("Leyendo archivo importado..."):
                    try:
                        if suffix == ".md":
                            summary_content = imported_file.read().decode("utf-8", errors="ignore")
                        elif suffix == ".pdf":
                            with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as temp_pdf:
                                temp_pdf.write(imported_file.read())
                                temp_pdf_path = temp_pdf.name
                            try:
                                reader = PdfReader(temp_pdf_path)
                                page_texts = []
                                for page in reader.pages:
                                    page_text = page.extract_text()
                                    if page_text:
                                        page_texts.append(page_text)
                                summary_content = "\n\n".join(page_texts)
                            finally:
                                if os.path.exists(temp_pdf_path):
                                    os.unlink(temp_pdf_path)
                        
                        if not summary_content.strip():
                            st.error("❌ El archivo importado no contiene texto legible.")
                        else:
                            save_class_data(
                                name=class_name_import.strip(),
                                summary=summary_content.strip(),
                                docs_extracted={},
                                depth="Importado desde archivo"
                            )
                            st.session_state.summary_output = summary_content.strip()
                            st.session_state.docs_extracted = {}
                            st.success(f"🎉 ¡Clase '{class_name_import.strip()}' importada con éxito y guardada en el historial!")
                            time.sleep(2)
                            st.rerun()
                    except Exception as e:
                        st.error(f"❌ Error al importar el archivo: {str(e)}")
        st.stop()

    # MODO CLÁSICO: PROCESAR DESDE GRABACIÓN
    st.markdown("### 🎙️ Cargar Nueva Sesión de Clase")
    
    class_name = st.text_input(
        "Nombre de la Clase", 
        value="", 
        placeholder="Ej: Clase 1 - Introduccion al Desarrollo Cognitivo",
        help="Escribe el nombre único con el que se guardará esta clase en tu base de datos local."
    )

    col1, col2 = st.columns(2)
    
    with col1:
        with st.container(border=True):
            st.markdown("#### 🎙️ Archivo de Audio o Video")
            media_file = st.file_uploader(
                "Sube la grabación de la clase", 
                type=["mp3", "wav", "mp4", "mov", "avi", "mkv", "m4a"],
                help="El sistema extraerá el audio de los archivos de video automáticamente. Soporta hasta 1000 MB."
            )
        
    with col2:
        with st.container(border=True):
            st.markdown("#### 📊 Materiales de Apoyo (Múltiples)")
            uploaded_docs = st.file_uploader(
                "Sube las diapositivas o documentos complementarios", 
                type=["pptx", "pdf"],
                accept_multiple_files=True,
                help="Puedes subir múltiples archivos de tipo .pptx y .pdf para complementar la explicación de la clase."
            )

    if media_file and uploaded_docs:
        st.markdown("<br>", unsafe_allow_html=True)
        
        if not api_key:
            st.error("⚠️ Debes proporcionar una Gemini API Key en el panel lateral.")
        elif not class_name.strip():
            st.warning("⚠️ Debes ingresar un nombre para la clase en el campo superior antes de procesar.")
        else:
            if st.button("✨ Procesar y Sincronizar Clase"):
                loader = show_full_screen_loader("Iniciando procesamiento de clase...")
                status_container = st.container()
                
                with status_container:
                    # 1. Guardar y procesar archivos locales
                    docs_extracted = {}
                    total_elements = 0
                    
                    for doc_file in uploaded_docs:
                        doc_name = doc_file.name
                        doc_suffix = Path(doc_name).suffix.lower()
                        is_pdf_doc = doc_suffix == ".pdf"
                        
                        update_full_screen_loader(loader, f"Extrayendo texto del archivo: {doc_name}...")
                        with tempfile.NamedTemporaryFile(delete=False, suffix=doc_suffix) as temp_doc:
                            temp_doc.write(doc_file.read())
                            temp_doc_path = temp_doc.name
                            
                            try:
                                if is_pdf_doc:
                                    extracted_parts = extract_text_from_pdf(temp_doc_path)
                                else:
                                    extracted_parts = extract_text_from_pptx(temp_doc_path)
                                docs_extracted[doc_name] = {
                                    "type": "PDF" if is_pdf_doc else "PPTX",
                                    "parts": extracted_parts
                                }
                                total_elements += len(extracted_parts)
                            except Exception as e:
                                st.error(f"❌ Error al procesar el archivo {doc_name}: {str(e)}")
                                loader.empty()
                                st.stop()
                            finally:
                                safe_remove_temp_file(temp_doc_path)
                                    
                    st.session_state.docs_extracted = docs_extracted
                    st.success(f"✅ Texto extraído con éxito de {len(uploaded_docs)} archivo(s) (total de {total_elements} páginas/diapositivas).")
                    
                    # 2. Manejar y optimizar el archivo de audio/video
                    media_suffix = Path(media_file.name).suffix.lower()
                    update_full_screen_loader(loader, "Procesando y extrayendo audio de tu grabación...")
                    with tempfile.NamedTemporaryFile(delete=False, suffix=media_suffix) as temp_media:
                        temp_media.write(media_file.read())
                        temp_media_path = temp_media.name
                    
                    audio_to_upload_path = None
                    is_video = media_suffix in [".mp4", ".mov", ".avi", ".mkv"]
                    
                    try:
                        if is_video:
                            update_full_screen_loader(loader, "Archivo de video detectado. Extrayendo canal de audio con moviepy...")
                            temp_audio_extracted = tempfile.NamedTemporaryFile(delete=False, suffix=".mp3")
                            audio_to_upload_path = temp_audio_extracted.name
                            temp_audio_extracted.close()
                            
                            extract_audio_from_video(temp_media_path, audio_to_upload_path)
                            st.success("✅ Audio extraído del video de forma correcta.")
                        else:
                            audio_to_upload_path = temp_media_path
                            st.success("✅ Archivo de audio listo para procesamiento.")
                            
                    except Exception as e:
                        st.error(f"❌ Error al procesar el audio/video: {str(e)}")
                        loader.empty()
                        st.stop()
                    finally:
                        if is_video and os.path.exists(temp_media_path):
                            os.unlink(temp_media_path)

                    # 3. Procesamiento y llamada a la API (Google o OpenAI)
                    try:
                        docs_payload = []
                        for filename, doc_data in docs_extracted.items():
                            doc_type = doc_data["type"]
                            parts = doc_data["parts"]
                            docs_payload.append(f"\n=========================================\n"
                                                f"CONTENIDO DEL MATERIAL DE APOYO: {filename} ({doc_type})\n"
                                                f"=========================================\n")
                            for p in parts:
                                unit_name = "PÁGINA" if doc_type == "PDF" else "DIAPOSITIVA"
                                part_info = f"--- {unit_name} NRO. {p['slide_num']} ---\n"
                                if p["content"]:
                                    part_info += f"[Texto]:\n{p['content']}\n"
                                if p["notes"]:
                                    part_info += f"[Notas del Orador]:\n{p['notes']}\n"
                                docs_payload.append(part_info)
                        
                        formatted_slides_text = "\n".join(docs_payload)
                        
                        # Determinar pautas de profundidad
                        if "Completo y Detallado" in summary_depth:
                            detail_instructions = """
- Genera un resumen académico de alta profundidad. Explaya cada concepto con al menos dos o tres párrafos detallados e incluye de forma textual los ejemplos clínicos o prácticos que mencione el profesor en el audio.
- No omitas ningún dato relevante, ejemplo práctico, analogía o concepto mencionado en la clase.
- El resumen debe ser un material de estudio completo que sirva para preparar un examen de alto nivel.
"""
                        elif "Estándar" in summary_depth:
                            detail_instructions = """
- Tu objetivo es generar un resumen estructurado, equilibrado y claro de la clase.
- Agrupa la explicación del docente en temas y secciones principales bien definidos.
- Resume las explicaciones verbales destacando los puntos más importantes de cada tema, sin extenderte en comentarios anecdóticos o redundancias, pero manteniendo el hilo lógico de la explicación.
"""
                        else:
                            detail_instructions = """
- Tu objetivo es generar un resumen conciso, sintético y muy acotado.
- Ve directo al grano. Utiliza listas con viñetas cortas, oraciones breves y resúmenes de un solo párrafo por tema principal.
- Identifica solo los conceptos fundamentales, eliminando detalles complementarios o explicaciones redundantes del audio.
"""

                        if "gpt-" in selected_model:
                            # --- FLUJO OPENAI: TRANSCRIPCIÓN CON WHISPER Y LUEGO CHAT COMPLETIONS ---
                            if not openai_api_key:
                                raise ValueError("Falta ingresar la OpenAI API Key en el panel lateral.")
                        
                            update_full_screen_loader(loader, "Transcribiendo grabación localmente con OpenAI Whisper...")
                            openai_client = OpenAI(api_key=openai_api_key)
                        
                            try:
                                audio_parts = split_audio_file_if_needed(audio_to_upload_path)
                            except Exception as e:
                                raise ValueError(f"Error al procesar/dividir el audio con MoviePy: {str(e)}")
                            
                            transcriptions = []
                            for idx, part_path in enumerate(audio_parts):
                                if len(audio_parts) > 1:
                                    update_full_screen_loader(loader, f"Transcribiendo parte {idx+1} de {len(audio_parts)} con OpenAI Whisper...")
                                with open(part_path, "rb") as audio_file:
                                    transcription_response = openai_client.audio.transcriptions.create(
                                        model="whisper-1",
                                        file=audio_file
                                    )
                                    transcriptions.append(transcription_response.text)
                            
                                if len(audio_parts) > 1 and part_path != audio_to_upload_path:
                                    try:
                                        os.unlink(part_path)
                                    except Exception:
                                        pass
                                    
                            transcription_text = " ".join(transcriptions)
                            st.success("✅ Archivo de audio transcrito correctamente mediante Whisper.")
                        
                            prompt = f"""
    Actúa como un profesor universitario y pedagogo experto. Te proporciono la transcripción completa del audio de la clase junto con el texto extraído de las diapositivas y documentos de apoyo que el docente usó como referencia.

    Tu tarea consiste en:
    1. Analizar minuciosamente la transcripción del audio (que representa la clase impartida real). El audio es la fuente principal de información y verdad para este resumen.
    2. Usar el texto de los documentos de apoyo únicamente como ayuda complementaria para enriquecer, completar o precisar el contenido del audio (por ejemplo, para obtener definiciones exactas de conceptos mencionados, ortografía de términos, o datos tabulados).
    3. Generar un resumen académico de la clase basado en las siguientes directrices de detalle:
    {detail_instructions}
    4. **CRÍTICO: No estructures el resumen diapositiva por diapositiva ni página por página**. En su lugar, el resumen debe fluir cronológicamente según los temas y secciones académicas desarrollados de forma verbal por el docente a lo largo de la clase.
    5. Al final del resumen, debes incluir un apartado titulado "📚 Glosario de Términos", donde definas detalladamente los conceptos clave, términos técnicos, fórmulas o datos más importantes explicados en la sesión.
    6. **CRÍTICO: EVITA dibujar tablas, diagramas o esquemas comparativos de dos columnas mediante espacios en blanco dentro de bloques de código (```). En su lugar, utiliza SIEMPRE tablas Markdown estándares (`| Columna 1 | Columna 2 |`) para cualquier estructura de columnas, comparaciones o alineación de términos.**

    Estructura tu reporte de la siguiente manera:

    # 🎓 Resumen Académico de la Clase (Nivel: {summary_depth.split(' ')[0]})

    ## 📌 Introducción y Contexto de la Sesión
    [Breve resumen de cómo inicia la clase, el tema general y los objetivos planteados por el docente.]

    ## 📝 Desarrollo Temático Detallado (Orden Cronológico de la Explicación por Audio)
    ### [Tema 1: Nombre del Tema]
    **Explicación y Conceptos:**
    [Desarrollo del tema aplicando las directrices de detalle seleccionadas. Sé coherente con la profundidad de análisis elegida.]
    **Ejemplos / Analogías / Casos Prácticos:**
    [Detalla cualquier ejemplo práctico, caso de estudio, anécdota o analogía mencionada en el audio para ilustrar el tema (si aplica al nivel de detalle seleccionado).]

    ### [Tema 2: Nombre del Tema]
    ... (repite la estructura para todos los temas principales tratados en el audio) ...

    ## 📚 Glosario de Términos
    * **[Término 1]**: [Definición y explicación de su relevancia en la clase].
    * **[Término 2]**: ...

    ---

    Aquí está la transcripción completa del audio de la clase:
    =========================================
    TRANSCRIPCIÓN DEL AUDIO:
    =========================================
    {transcription_text}

    Aquí está el texto extraído de los documentos de apoyo para tu referencia y complemento:
    =========================================
    DOCUMENTOS DE APOYO:
    =========================================
    {formatted_slides_text}

    Por favor, elabora el resumen en idioma español, asegurándose de mantener un alto estándar académico y formal.
    """
                            update_full_screen_loader(loader, f"Generando el resumen estructurado con OpenAI ({selected_model})...")
                            response = openai_client.chat.completions.create(
                                model=selected_model,
                                messages=[{"role": "user", "content": prompt}]
                            )
                            response_text = response.choices[0].message.content
                        else:
                            # --- FLUJO GEMINI: SUBIR AUDIO POR FILE API ---
                            if not api_key:
                                raise ValueError("Falta ingresar la Gemini API Key en el panel lateral.")
                        
                            client = genai.Client(api_key=api_key)
                        
                            docs_payload = []
                            for filename, doc_data in docs_extracted.items():
                                doc_type = doc_data["type"]
                                parts = doc_data["parts"]
                                docs_payload.append(f"\n=========================================\n"
                                                    f"CONTENIDO DEL MATERIAL DE APOYO: {filename} ({doc_type})\n"
                                                    f"=========================================\n")
                                for p in parts:
                                    unit_name = "PÁGINA" if doc_type == "PDF" else "DIAPOSITIVA"
                                    part_info = f"--- {unit_name} NRO. {p['slide_num']} ---\n"
                                    if p["content"]:
                                        part_info += f"[Texto]:\n{p['content']}\n"
                                    if p["notes"]:
                                        part_info += f"[Notas del Orador]:\n{p['notes']}\n"
                                    docs_payload.append(part_info)
                        
                            formatted_slides_text = "\n".join(docs_payload)
                        
                            update_full_screen_loader(loader, "Subiendo grabación a los servidores de Gemini...")
                            uploaded_file = client.files.upload(file=audio_to_upload_path)
                            
                            while uploaded_file.state.name == "PROCESSING":
                                update_full_screen_loader(loader, "La grabación se está procesando en los servidores de Gemini, por favor espera...")
                                time.sleep(3)
                                uploaded_file = client.files.get(name=uploaded_file.name)
                        
                            if uploaded_file.state.name == "FAILED":
                                raise ValueError("Gemini falló al procesar el archivo de audio subido.")
                            
                            st.success("✅ Archivo procesado correctamente en la API de Gemini.")
                        
                            prompt = f"""
    Actúa como un profesor universitario y pedagogo experto. He subido el audio de una clase y te adjunto a continuación el contenido de texto extraído de las diapositivas y documentos de apoyo que el docente usó como referencia.

    Tu tarea consiste en:
    1. Analizar minuciosamente la explicación del audio (que representa la clase impartida real). El audio es la fuente principal de información y verdad para este resumen.
    2. Usar el texto de los documentos de apoyo únicamente como ayuda complementaria para enriquecer, completar o precisar el contenido del audio (por ejemplo, para obtener definiciones exactas de conceptos mencionados, ortografía de términos, o datos tabulados).
    3. Generar un resumen académico de la clase basado en las siguientes directrices de detalle:
    {detail_instructions}
    4. **CRÍTICO: No estructures el resumen diapositiva por diapositiva ni página por página**. En su lugar, el resumen debe fluir cronológicamente según los temas y secciones académicas desarrollados de forma verbal por el docente a lo largo de la clase.
    5. Al final del resumen, debes incluir un apartado titulado "📚 Glosario de Términos", donde definas detalladamente los conceptos clave, términos técnicos, fórmulas o datos más importantes explicados en la sesión.
    6. **CRÍTICO: EVITA dibujar tablas, diagramas o esquemas comparativos de dos columnas mediante espacios en blanco dentro de bloques de código (```). En su lugar, utiliza SIEMPRE tablas Markdown estándares (`| Columna 1 | Columna 2 |`) para cualquier estructura de columnas, comparaciones o alineación de términos.**

    Estructura tu reporte de la siguiente manera:

    # 🎓 Resumen Académico de la Clase (Nivel: {summary_depth.split(' ')[0]})

    ## 📌 Introducción y Contexto de la Sesión
    [Breve resumen de cómo inicia la clase, el tema general y los objetivos planteados por el docente.]

    ## 📝 Desarrollo Temático Detallado (Orden Cronológico de la Explicación por Audio)
    ### [Tema 1: Nombre del Tema]
    **Explicación y Conceptos:**
    [Desarrollo del tema aplicando las directrices de detalle seleccionadas. Sé coherente con la profundidad de análisis elegida.]
    **Ejemplos / Analogías / Casos Prácticos:**
    [Detalla cualquier ejemplo práctico, caso de estudio, anécdota o analogía mencionada en el audio para ilustrar el tema (si aplica al nivel de detalle seleccionado).]

    ### [Tema 2: Nombre del Tema]
    ... (repite la estructura para todos los temas principales tratados en el audio) ...

    ## 📚 Glosario de Términos
    * **[Término 1]**: [Definición y explicación de su relevancia en la clase].
    * **[Término 2]**: ...

    ---

    Aquí está el texto extraído de los documentos de apoyo para tu referencia y complemento:
    {formatted_slides_text}

    Por favor, elabora el resumen en idioma español, asegurándose de mantener un alto estándar académico y formal.
    """
                            update_full_screen_loader(loader, "Analizando grabación y generando el resumen estructurado con Gemini...")
                            response = client.models.generate_content(
                                model=selected_model,
                                contents=[uploaded_file, prompt]
                            )
                            response_text = response.text

                        save_class_data(
                            name=class_name.strip(),
                            summary=response_text,
                            docs_extracted=docs_extracted,
                            depth=summary_depth
                        )
                        st.success(f"🎉 ¡Análisis generado con éxito y guardado en el historial como '{class_name.strip()}'!")
                        
                    except Exception as e:
                        st.error(handle_api_error(e))
                    finally:
                        loader.empty()
                        
        if st.session_state.summary_output:
            st.markdown("---")
            st.markdown(st.session_state.summary_output)
            
            # Botones de descarga para el resumen generado
            pdf_bytes = None
            try:
                pdf_bytes = generate_pdf_from_markdown(st.session_state.summary_output)
            except Exception as e:
                st.error(f"Error al generar PDF: {str(e)}")
                
            col_down_md, col_down_pdf = st.columns(2)
            with col_down_md:
                st.download_button(
                    label="📥 Descargar Resumen (Markdown)",
                    data=st.session_state.summary_output,
                    file_name=f"{output_filename}.md",
                    mime="text/markdown",
                    key="tab1_md_btn",
                    use_container_width=True
                )
            with col_down_pdf:
                if pdf_bytes:
                    st.download_button(
                        label="📥 Descargar Resumen (PDF)",
                        data=pdf_bytes,
                        file_name=f"{output_filename}.pdf",
                        mime="application/pdf",
                        key="tab1_pdf_btn",
                        use_container_width=True
                    )
            
            # Mostrar material de apoyo extraído
            st.markdown("---")
            st.markdown("### 🖥️ Material de Apoyo Extraído")
            st.markdown("A continuación se muestra el contenido extraído de los documentos de apoyo para esta clase.")
            
            for filename, doc_data in st.session_state.docs_extracted.items():
                st.markdown(f"#### 📄 Archivo: `{filename}` ({doc_data['type']})")
                parts = doc_data["parts"]
                unit_name = "Página" if doc_data["type"] == "PDF" else "Diapositiva"
                for s in parts:
                    with st.expander(f"🖥️ {unit_name} {s['slide_num']}"):
                        if s["content"]:
                            st.code(s["content"], language="text")
                        else:
                            st.info("Sin texto detectado.")
                        if doc_data["type"] == "PPTX" and s["notes"]:
                            st.markdown("**Notas del Orador:**")
                            st.write(s["notes"])

# =====================================================================
# PESTAÑA 2: HISTORIAL DE CLASES
# =====================================================================
with tab_history:
    st.markdown("### 📚 Historial de Clases Procesadas")
    saved_classes_list = list_saved_classes()
    
    if not saved_classes_list:
        st.info("📂 Aún no tienes clases guardadas. Procesa tu primera clase en la pestaña 🎙️ **Procesar Nueva Clase**.")
    else:
        selected_saved_class = st.selectbox(
            "Selecciona una clase guardada para visualizar",
            options=saved_classes_list
        )
        
        class_data = load_class_data(selected_saved_class)
        
        if class_data:
            st.markdown(f"## 📂 Clase: `{class_data['name']}`")
            st.caption(f"📅 **Fecha de guardado**: {class_data.get('date', 'Desconocida')} | ⚙️ **Nivel original**: {class_data.get('depth', 'Desconocido')}")
            st.markdown("---")
            
            tab_summary_hist, tab_original_hist = st.tabs(["📝 Resumen de Clase", "📂 Texto de Apoyo Original"])
            
            with tab_summary_hist:
                st.markdown(class_data["summary"])
                
                try:
                    pdf_bytes_hist = generate_pdf_from_markdown(class_data["summary"])
                except Exception as e:
                    pdf_bytes_hist = None
                    st.error(f"Error al generar el PDF del historial: {str(e)}")
                    
                col_hist_md, col_hist_pdf = st.columns(2)
                with col_hist_md:
                    st.download_button(
                        label="📥 Descargar Resumen (Markdown)",
                        data=class_data["summary"],
                        file_name=f"{selected_saved_class}.md",
                        mime="text/markdown",
                        key="hist_md_btn",
                        use_container_width=True
                    )
                with col_hist_pdf:
                    if pdf_bytes_hist:
                        st.download_button(
                            label="📥 Descargar Resumen (PDF)",
                            data=pdf_bytes_hist,
                            file_name=f"{selected_saved_class}.pdf",
                            mime="application/pdf",
                            key="hist_pdf_btn",
                            use_container_width=True
                        )
                
                st.markdown("---")
                st.markdown("### 🔄 Regenerar con Otro Nivel de Detalle")
                st.markdown("*(Esta acción no requiere volver a subir el archivo de audio o video)*")
                
                new_depth = st.selectbox(
                    "Elige el nuevo nivel de profundidad para esta clase",
                    options=[
                        "Completo y Detallado (Estudio exhaustivo)",
                        "Estándar (Equilibrado)",
                        "Acotado (Conciso y directo)"
                    ],
                    index=0,
                    key="history_depth_select"
                )
                
                if st.button("🔄 Cambiar Profundidad y Actualizar"):
                    if not api_key and "gpt-" not in selected_model:
                        st.error("🔑 Ingresa una Gemini API Key en el panel lateral para continuar.")
                    elif "gpt-" in selected_model and not openai_api_key:
                        st.error("🔑 Ingresa una OpenAI API Key en el panel lateral para continuar.")
                    else:
                        loader = show_full_screen_loader("Regenerando resumen con el nuevo nivel de detalle...")
                        try:
                            if "Completo y Detallado" in new_depth:
                                new_instructions = "Genera un resumen académico EXHAUSTIVO, extenso y sumamente detallado de la clase. No omitas ningún concepto o explicación importante."
                            elif "Estándar" in new_depth:
                                new_instructions = "Genera un resumen académico estructurado, equilibrado y claro, agrupando los conceptos en secciones principales."
                            else:
                                new_instructions = "Genera un resumen académico conciso, sintético y acotado, directo al grano utilizando listas y explicaciones cortas."
                            
                            prompt_regen = f"""
Actúa como un profesor universitario y pedagogo experto. Te proporciono el resumen actual de una clase y el contenido textual de los materiales de apoyo de la misma.

Tu tarea consiste en:
1. Regenerar el resumen adaptándolo estrictamente al nuevo nivel de profundidad seleccionado: {new_depth}.
2. Seguir las directrices de detalle: {new_instructions}
3. Mantener la coherencia académica y asegurar que el resumen fluya lógicamente de principio a fin, terminando con un glosario de términos.
4. **CRÍTICO: EVITA dibujar tablas, diagramas o esquemas comparativos de dos columnas mediante espacios en blanco dentro de bloques de código (```). En su lugar, utiliza SIEMPRE tablas Markdown estándares (`| Columna 1 | Columna 2 |`) para cualquier estructura de columnas, comparaciones o alineación de términos.**

### RESUMEN ANTERIOR DE LA CLASE:
{class_data['summary']}

---
Por favor, genera el nuevo resumen adaptado en español de forma limpia y formateada.
"""
                            if "gpt-" in selected_model:
                                openai_client = OpenAI(api_key=openai_api_key)
                                response_regen = openai_client.chat.completions.create(
                                    model=selected_model,
                                    messages=[{"role": "user", "content": prompt_regen}]
                                )
                                response_regen_text = response_regen.choices[0].message.content
                            else:
                                client = genai.Client(api_key=api_key)
                                response_regen = client.models.generate_content(
                                    model=selected_model,
                                    contents=[prompt_regen]
                                )
                                response_regen_text = response_regen.text
                                
                            save_class_data(
                                name=selected_saved_class,
                                summary=response_regen_text,
                                docs_extracted=class_data["docs_extracted"],
                                depth=new_depth
                            )
                            st.success("🎉 ¡El resumen de la clase ha sido actualizado con éxito!")
                            st.rerun()
                        except Exception as e:
                            st.error(handle_api_error(e))
                        finally:
                            loader.empty()
                                
            with tab_original_hist:
                if class_data.get("docs_extracted"):
                    for filename, doc_data in class_data["docs_extracted"].items():
                        st.markdown(f"#### 📄 Archivo: `{filename}` ({doc_data['type']})")
                        parts = doc_data["parts"]
                        unit_name = "Página" if doc_data["type"] == "PDF" else "Diapositiva"
                        for s in parts:
                            with st.expander(f"🖥️ {unit_name} {s['slide_num']}"):
                                if s["content"]:
                                    st.code(s["content"], language="text")
                                else:
                                    st.info("Sin texto detectado.")
                                if doc_data["type"] == "PPTX" and s["notes"]:
                                    st.markdown("**Notas del Orador:**")
                                    st.write(s["notes"])
                else:
                    st.info("No hay materiales de texto guardados en esta clase.")
                    
        # --- Panel de Administración ---
        st.markdown("---")
        with st.expander("⚙️ Administrar Clases Guardadas (Renombrar / Eliminar)"):
            st.markdown("#### ⚙️ Panel de Administración")
            
            admin_class_select = st.selectbox(
                "Selecciona una clase para renombrar o eliminar",
                options=saved_classes_list,
                key="admin_class_select"
            )
            
            col_admin_rename, col_admin_delete = st.columns(2)
            
            with col_admin_rename:
                st.markdown("##### ✏️ Renombrar Clase")
                new_name_input = st.text_input("Nuevo nombre de la clase", value=admin_class_select)
                if st.button("✏️ Renombrar clase", key="admin_rename_btn"):
                    if not new_name_input.strip():
                        st.warning("⚠️ Ingresa un nombre válido.")
                    else:
                        try:
                            old_path = get_classes_dir() / f"{admin_class_select}.json"
                            new_path = get_classes_dir() / f"{new_name_input.strip()}.json"
                            if new_path.exists():
                                st.error("❌ Ya existe una clase con ese nombre.")
                            else:
                                data = load_class_data(admin_class_select)
                                if data:
                                    data["name"] = new_name_input.strip()
                                    with open(new_path, 'w', encoding='utf-8') as f:
                                        json.dump(data, f, ensure_ascii=False, indent=4)
                                    os.unlink(old_path)
                                    st.success(f"✏️ Clase renombrada con éxito a '{new_name_input.strip()}'")
                                    time.sleep(2)
                                    st.rerun()
                        except Exception as e:
                            st.error(f"❌ Error al renombrar la clase: {str(e)}")
                            
            with col_admin_delete:
                st.markdown("##### 🗑️ Eliminar Clase")
                st.write(f"¿Estás seguro de que deseas eliminar permanentemente la clase `{admin_class_select}`?")
                st.warning("⚠️ Esta acción es irreversible y eliminará el resumen e historial de la clase.")
                if st.button("🗑️ Eliminar permanentemente", key="admin_delete_btn"):
                    try:
                        file_path = get_classes_dir() / f"{admin_class_select}.json"
                        if file_path.exists():
                            os.unlink(file_path)
                            st.success(f"🗑️ Clase '{admin_class_select}' eliminada con éxito.")
                            time.sleep(2)
                            st.rerun()
                    except Exception as e:
                        st.error(f"❌ Error al eliminar la clase: {str(e)}")

# =====================================================================
# PESTAÑA 3: ZONA DE ESTUDIO E INTEGRACIÓN
# =====================================================================
with tab_study:
    st.markdown('<h1 class="main-title">🎯 Zona de Estudio e Integración</h1>', unsafe_allow_html=True)
    st.markdown("<p class='subtitle-text'>Genera exámenes simulados, cuestionarios y flashcards individuales o consolidadas integrando tus clases y lecturas.</p>", unsafe_allow_html=True)
    
    col_single, col_multi = st.columns(2)
    
    with col_single:
        st.markdown("### 🏫 Estudio Individual")
        st.markdown("Genera material de repaso a partir del resumen de una clase en particular.")
        
        classes_list = list_saved_classes()
        if not classes_list:
            st.warning("💡 Primero debes procesar y guardar al menos una clase en el historial.")
        else:
            selected_single_class = st.selectbox(
                "Selecciona una clase",
                options=classes_list,
                key="study_single_class_select"
            )
            
            study_tool = st.selectbox(
                "Elige la herramienta de estudio",
                options=["Cuestionario de Práctica (Quiz)", "Fichas de Repaso (Flashcards)"]
            )
            
            # Selector de cantidad de preguntas/fichas
            study_qty = st.slider(
                f"Cantidad de {study_tool.split(' ')[0]} a generar",
                min_value=5,
                max_value=30,
                value=10,
                step=5,
                key="study_qty_slider"
            )
            
            if st.button("✨ Generar Material de Estudio", key="btn_gen_single_study"):
                if not api_key and "gpt-" not in selected_model:
                    st.error("🔑 Ingresa tu Gemini API Key en el panel lateral.")
                elif "gpt-" in selected_model and not openai_api_key:
                    st.error("🔑 Ingresa tu OpenAI API Key en el panel lateral.")
                else:
                    loader = show_full_screen_loader(f"Generando {study_tool} con IA...")
                    try:
                        class_data = load_class_data(selected_single_class)
                        source_content = class_data["summary"]
                        
                        if study_tool == "Cuestionario de Práctica (Quiz)":
                            prompt_study = f"""
Actúa como un profesor universitario. Genera un cuestionario de práctica (quiz) con {study_qty} preguntas de opción múltiple en español basadas estrictamente en la siguiente fuente de estudio.
Para cada pregunta, proporciona 4 opciones (A, B, C, D) y la explicación de la respuesta correcta.

### DIRECTRICES DE FORMATO DE OPCIONES (MUY IMPORTANTE):
- Presenta cada opción en su propia línea, con un espacio libre antes y después (doble salto de línea).
- Agrega un emoji azul `🔹 ` al inicio de cada opción.
- Formatea el indicador de la opción en negrita, seguido de un paréntesis y el texto. Ejemplo:
  🔹 **A)** [Texto de la opción]
  
  🔹 **B)** [Texto de la opción]
  
  🔹 **C)** [Texto de la opción]
  
  🔹 **D)** [Texto de la opción]
- No dejes las opciones juntas ni en la misma línea. Asegura un espaciado amplio y limpio para máxima legibilidad.

Oculta la respuesta y explicación en un bloque desplegable Markdown formateado así:

### Pregunta X: [Enunciado]

🔹 **A)** [Opción A]

🔹 **B)** [Opción B]

🔹 **C)** [Opción C]

🔹 **D)** [Opción D]

:::spoiler Revelar Respuesta Correcta
**Respuesta Correcta:** [Indicar letra]
**Explicación:** [Justificación detallada del porqué es correcta...]
:::

### FUENTE DE ESTUDIO:
{source_content}
"""
                        else: # Flashcards
                            prompt_study = f"""
Actúa como un tutor de memorización. Genera un conjunto de {study_qty} fichas de repaso (flashcards) que abarquen los conceptos más críticos y difíciles de la siguiente fuente de estudio.
Oculta la respuesta en un bloque desplegable Markdown formateado así:

### Ficha X: [Término o Pregunta Clave]
:::spoiler Voltear Tarjeta (Ver Definición/Respuesta)
**Explicación:** [Explicación detallada y ejemplos...]
:::

### FUENTE DE ESTUDIO:
{source_content}
"""
                        if "gpt-" in selected_model:
                            openai_client = OpenAI(api_key=openai_api_key)
                            response_study = openai_client.chat.completions.create(
                                model=selected_model,
                                messages=[{"role": "user", "content": prompt_study}]
                            )
                            response_text = response_study.choices[0].message.content
                        else:
                            client = genai.Client(api_key=api_key)
                            response_study = client.models.generate_content(
                                model=selected_model,
                                contents=[prompt_study]
                            )
                            response_text = response_study.text
                            
                        st.session_state.study_material = response_text
                        st.session_state.study_title = f"{study_tool} - {selected_single_class}"
                        st.session_state.consolidated_material = None
                        st.success("🎉 ¡Material de estudio generado correctamente!")
                    except Exception as e:
                        st.error(handle_api_error(e))
                    finally:
                        loader.empty()
                        
    with col_multi:
        st.markdown("### 📚 Integración de Fuentes Múltiples")
        st.markdown("Consolida y unifica el contenido de múltiples clases o lecturas de apoyo en un único material.")
        
        active_sub = get_active_subject()
        classes_list = list_saved_classes()
        readings_list = list_saved_readings()
        
        all_sources = [f"🏫 Clase: {c}" for c in classes_list] + [f"📖 Lectura: {r}" for r in readings_list]
        
        if len(all_sources) < 2:
            st.warning("💡 Necesitas tener al menos 2 fuentes guardadas (clases o lecturas) en el historial para realizar una integración de fuentes múltiples.")
        else:
            selected_mix = st.multiselect(
                "Selecciona dos o más fuentes del historial",
                options=all_sources,
                help="Elige las clases y/o lecturas de apoyo de las cuales deseas crear una guía de estudio o examen conjunto."
            )
            
            consolidate_action = st.selectbox(
                "Elige qué material unificado generar",
                options=[
                    "Guía de Estudio Consolidada (Temario unificado)",
                    "Examen de Práctica Consolidado (Quiz Integrado)",
                    "Fichas de Repaso Consolidadas"
                ]
            )
            
            multi_quiz_qty = 15
            multi_flash_qty = 15
            guide_depth = "Intermedia"
            
            if "Guía de Estudio" in consolidate_action:
                guide_depth = st.radio(
                    "Profundidad de la Guía de Estudio",
                    options=["Completa y Detallada", "Intermedia", "Acotada"],
                    horizontal=True,
                    index=1,
                    key="multi_guide_depth_radio",
                    help="Completa: Muy detallada y extensa. Intermedia: Estructura estándar equilibrada. Acotada: Síntesis concisa."
                )
            elif "Examen" in consolidate_action:
                multi_quiz_qty = st.slider(
                    "Cantidad de preguntas unificadas",
                    min_value=5,
                    max_value=40,
                    value=15,
                    step=5,
                    key="multi_quiz_qty_slider"
                )
            else: # Flashcards
                multi_flash_qty = st.slider(
                    "Cantidad de fichas consolidadas",
                    min_value=5,
                    max_value=40,
                    value=15,
                    step=5,
                    key="multi_flash_qty_slider"
                )
            
            if st.button("✨ Generar Material Consolidado", key="btn_gen_multi_study"):
                if len(selected_mix) < 2:
                    st.error("⚠️ Debes seleccionar al menos 2 fuentes para consolidar.")
                elif not api_key and "gpt-" not in selected_model:
                    st.error("🔑 Ingresa tu Gemini API Key en el panel lateral.")
                elif "gpt-" in selected_model and not openai_api_key:
                    st.error("🔑 Ingresa tu OpenAI API Key en el panel lateral.")
                else:
                    loader = show_full_screen_loader(f"Generando material consolidado ({consolidate_action})...")
                    try:
                        compiled_texts = []
                        for src in selected_mix:
                            if src.startswith("🏫 Clase: "):
                                c_name = src[len("🏫 Clase: "):]
                                c_data = load_class_data(c_name)
                                if c_data:
                                    compiled_texts.append(f"### RESUMEN DE CLASE: {c_data['name']}\n{c_data['summary']}")
                            elif src.startswith("📖 Lectura: "):
                                r_name = src[len("📖 Lectura: "):]
                                r_data = load_reading_data(r_name)
                                if r_data:
                                    r_text = get_reading_consolidated_text(r_data)
                                    compiled_texts.append(f"### LECTURA DE APOYO: {r_data['name']}\n{r_text}")
                        
                        all_sources_text = "\n\n".join(compiled_texts)
                        
                        if "Guía de Estudio" in consolidate_action:
                            if guide_depth == "Completa y Detallada":
                                depth_instruction = "Genera una guía de estudio extremadamente exhaustiva y de extensión completa, desglosando minuciosamente cada tema, subtema, definición, ejemplo y analogía sin recortar nada de la materia académica."
                            elif guide_depth == "Acotada":
                                depth_instruction = "Genera una guía de estudio acotada y muy sintetizada, resumiendo los conceptos clave en oraciones breves, explicaciones directas y viñetas sintéticas."
                            else: # Intermedia
                                depth_instruction = "Genera una guía de estudio de profundidad intermedia, balanceando el detalle conceptual con una estructura clara y concisa de los temas principales."
                            
                            prompt_multi = f"""
Actúa como un profesor universitario y director de cátedra. Te proporciono el contenido de varias fuentes de estudio (que incluyen clases y lecturas).
Tu tarea es unificar e integrar todo este contenido en una única "Guía de Estudio Consolidada" adaptada al siguiente nivel de detalle:
- Nivel de detalle: {guide_depth}
- Directriz: {depth_instruction}

La guía debe:
1. Tener una estructura coherente y unificada, agrupando los temas similares.
2. Evitar redundancias pero retener todo el nivel de detalle académico importante correspondiente a la profundidad solicitada.
3. Incluir introducciones de transición que conecten los conceptos de las diferentes fuentes.
4. Finalizar con un "Glosario Maestro" de términos unificado.
5. **CRÍTICO: EVITA dibujar tablas, diagramas o esquemas comparativos de dos columnas mediante espacios en blanco dentro de bloques de código (```). En su lugar, utiliza SIEMPRE tablas Markdown estándares (`| Columna 1 | Columna 2 |`) para cualquier estructura de columnas, comparaciones o alineación de términos.**

### FUENTES DE ESTUDIO A CONSOLIDAR:
{all_sources_text}
"""
                        elif "Examen" in consolidate_action:
                            prompt_multi = f"""
Actúa como un evaluador académico. Genera un examen simulado unificado con {multi_quiz_qty} preguntas de opción múltiple en español basadas en las fuentes de estudio provistas (resúmenes de clases y lecturas).
Las preguntas deben cubrir equitativamente los temas de todas las fuentes seleccionadas.
Para cada pregunta, proporciona 4 opciones (A, B, C, D) y la explicación de la respuesta correcta.

### DIRECTRICES DE FORMATO DE OPCIONES (MUY IMPORTANTE):
- Presenta cada opción en su propia línea, con un espacio libre antes y después (doble salto de línea).
- Agrega un emoji azul `🔹 ` al inicio de cada opción.
- Formatea el indicador de la opción en negrita, seguido de un paréntesis y el texto. Ejemplo:
  🔹 **A)** [Texto de la opción]
  
  🔹 **B)** [Texto de la opción]
  
  🔹 **C)** [Texto de la opción]
  
  🔹 **D)** [Texto de la opción]
- No dejes las opciones juntas ni en la misma línea. Asegura un espaciado amplio y limpio para máxima legibilidad.

Oculta la respuesta y explicación en un bloque desplegable Markdown formateado así:

### Pregunta X: [Enunciado]

🔹 **A)** [Opción A]

🔹 **B)** [Opción B]

🔹 **C)** [Opción C]

🔹 **D)** [Opción D]

:::spoiler Revelar Respuesta Correcta
**Respuesta Correcta:** [Indicar letra]
**Explicación:** [Justificación detallada del porqué es correcta...]
:::

### FUENTES DE ESTUDIO A EVALUAR:
{all_sources_text}
"""
                        else: # Fichas de Repaso
                            prompt_multi = f"""
Actúa como un tutor de memorización. Genera un conjunto de {multi_flash_qty} fichas de repaso (flashcards) consolidadas que abarquen los conceptos más críticos y difíciles de todas las fuentes de estudio seleccionadas.
Oculta la respuesta en un bloque desplegable Markdown formateado así:

### Ficha X: [Término o Pregunta Clave]
:::spoiler Voltear Tarjeta (Ver Definición/Respuesta)
**Explicación:** [Explicación detallada y ejemplos...]
:::

### FUENTES DE ESTUDIO A REPASAR:
{all_sources_text}
"""
                        if "gpt-" in selected_model:
                            openai_client = OpenAI(api_key=openai_api_key)
                            response_multi = openai_client.chat.completions.create(
                                model=selected_model,
                                messages=[{"role": "user", "content": prompt_multi}]
                            )
                            response_text = response_multi.choices[0].message.content
                        else:
                            client = genai.Client(api_key=api_key)
                            response_multi = client.models.generate_content(
                                model=selected_model,
                                contents=[prompt_multi]
                            )
                            response_text = response_multi.text
                            
                        st.session_state.consolidated_material = response_text
                        st.session_state.consolidated_title = f"{consolidate_action} (Consolidado)"
                        st.session_state.study_material = None
                        st.success("🎉 ¡Material consolidado generado correctamente!")
                    except Exception as e:
                        st.error(handle_api_error(e))
                    finally:
                        loader.empty()

    # RENDERIZAR RESULTADO DE ESTUDIO INDIVIDUAL
    if st.session_state.study_material:
        st.markdown("---")
        st.markdown(f"## 📖 {st.session_state.study_title}")
        render_study_material_with_expanders(st.session_state.study_material)
        
        download_text_study = st.session_state.study_material
        if "Quiz" in st.session_state.study_title:
            download_text_study = adjust_quiz_answers_for_download(st.session_state.study_material)
        
        st.markdown("#### 📥 Descargar Material de Estudio")
        study_filename = st.text_input(
            "Nombre de archivo de descarga",
            value=f"estudio_{st.session_state.study_title.lower().replace(' ', '_').replace(':', '')}",
            key="study_filename_input"
        )
        
        pdf_bytes_study = None
        try:
            pdf_bytes_study = generate_pdf_from_markdown(download_text_study)
        except Exception as e:
            st.error(f"Error al generar PDF: {str(e)}")
            
        col_study_md, col_study_pdf = st.columns(2)
        with col_study_md:
            st.download_button(
                label="📥 Descargar Material (Markdown)",
                data=download_text_study,
                file_name=f"{study_filename}.md",
                mime="text/markdown",
                key="study_md_btn",
                use_container_width=True
            )
        with col_study_pdf:
            if pdf_bytes_study:
                st.download_button(
                    label="📥 Descargar Material (PDF)",
                    data=pdf_bytes_study,
                    file_name=f"{study_filename}.pdf",
                    mime="application/pdf",
                    key="study_pdf_btn",
                    use_container_width=True
                )

    # RENDERIZAR RESULTADO DE ESTUDIO CONSOLIDADO
    if st.session_state.consolidated_material:
        st.markdown("---")
        st.markdown(f"## 📖 {st.session_state.consolidated_title}")
        st.markdown(st.session_state.consolidated_material)
        
        download_text_consolidated = st.session_state.consolidated_material
        if "Quiz" in st.session_state.consolidated_title or "Examen" in st.session_state.consolidated_title:
            download_text_consolidated = adjust_quiz_answers_for_download(st.session_state.consolidated_material)
            
        st.markdown("#### 📥 Descargar Material Consolidado")
        consolidated_filename = st.text_input(
            "Nombre de archivo de descarga",
            value=f"consolidado_{st.session_state.consolidated_title.lower().replace(' ', '_').replace(':', '')}",
            key="consolidated_filename_input"
        )
        
        pdf_bytes_consolidated = None
        try:
            pdf_bytes_consolidated = generate_pdf_from_markdown(download_text_consolidated)
        except Exception as e:
            st.error(f"Error al generar PDF: {str(e)}")
            
        col_con_md, col_con_pdf = st.columns(2)
        with col_con_md:
            st.download_button(
                label="📥 Descargar Material Consolidado (Markdown)",
                data=download_text_consolidated,
                file_name=f"{consolidated_filename}.md",
                mime="text/markdown",
                key="con_md_btn",
                use_container_width=True
            )
        with col_con_pdf:
            if pdf_bytes_consolidated:
                st.download_button(
                    label="📥 Descargar Material Consolidado (PDF)",
                    data=pdf_bytes_consolidated,
                    file_name=f"{consolidated_filename}.pdf",
                    mime="application/pdf",
                    key="con_pdf_btn",
                    use_container_width=True
                )

# =====================================================================
# PESTAÑA 4: LECTURAS DE APOYO
# =====================================================================
with tab_readings:
    st.markdown('<h1 class="main-title">📖 Lecturas de Apoyo</h1>', unsafe_allow_html=True)
    st.markdown("<p class='subtitle-text'>Sube libros, capítulos o apuntes complementarios para generar resúmenes, realizar OCR por IA, o extraer capítulos específicos.</p>", unsafe_allow_html=True)
    
    active_sub = get_active_subject()
    
    col_upload_reading, col_saved_readings = st.columns(2)
    
    with col_upload_reading:
        st.markdown("### 📤 Subir Lectura de Apoyo")
        
        uploaded_reading_file = st.file_uploader(
            "Sube un libro, capítulo o material complementario para estudiar (PDF o PPTX)",
            type=["pdf", "pptx"],
            key="uploaded_reading_file"
        )
        
        if uploaded_reading_file:
            st.markdown("#### 👁️ Previsualizar y Leer Archivo Subido (Vista Rápida)")
            
            temp_filename = f"temp_preview_{uploaded_reading_file.name}"
            with open(temp_filename, "wb") as f:
                f.write(uploaded_reading_file.getbuffer())
                
            preview_pages = []
            suffix = Path(uploaded_reading_file.name).suffix.lower()
            try:
                preview_pages = extract_reading_pages(temp_filename, suffix)
            except Exception as e:
                st.error(f"Error al analizar previsualización: {str(e)}")
                
            if os.path.exists(temp_filename):
                os.remove(temp_filename)
                
            if not preview_pages:
                st.warning("⚠️ No se pudo extraer texto legible para la previsualización.")
            else:
                preview_page_num = st.number_input(
                    f"Página a previsualizar (1 a {len(preview_pages)})",
                    min_value=1,
                    max_value=len(preview_pages),
                    value=1,
                    key="preview_page_select"
                )
                
                selected_page_text = preview_pages[preview_page_num - 1]["content"]
                st.markdown(f"**Vista Previa de la Página {preview_page_num}**")
                
                # Inicializar cache de vista previa OCR
                if "ocr_preview_cache" not in st.session_state:
                    st.session_state.ocr_preview_cache = {}
                    
                cache_key = f"{uploaded_reading_file.name}_p{preview_page_num}"
                
                if not selected_page_text.strip() and cache_key not in st.session_state.ocr_preview_cache:
                    st.warning("⚠️ Esta página parece estar vacía (posiblemente un escaneo o imagen).")
                    if st.button("🔍 Transcribir esta página con IA (OCR)", key=f"btn_ocr_p_{preview_page_num}"):
                        if not api_key and "gpt-" not in selected_model:
                            st.error("🔑 Ingresa tu Gemini API Key en el panel lateral.")
                        elif "gpt-" in selected_model and not openai_api_key:
                            st.error("🔑 Ingresa tu OpenAI API Key en el panel lateral.")
                        else:
                            with st.spinner("Transcribiendo página..."):
                                try:
                                    # Guardar de nuevo temporalmente para OCR
                                    temp_ocr_path = f"temp_ocr_{uploaded_reading_file.name}"
                                    with open(temp_ocr_path, "wb") as f:
                                        f.write(uploaded_reading_file.getbuffer())
                                        
                                    if "gpt-" in selected_model:
                                        # Para OpenAI no subimos el archivo directamente a Whisper si supera límite, 
                                        # pero para OCR de una página podemos usar GPT-4V/GPT-4o enviando la imagen de la página.
                                        # Sin embargo, como el PDF está guardado en disco, podemos usar el cliente de Gemini que soporta PDF nativo.
                                        # Así que de forma segura usamos Gemini para OCR de PDF si está disponible, o mostramos error.
                                        st.info("💡 El OCR directo de PDFs se realiza a través de los servidores de Gemini API.")
                                        ocr_client = genai.Client(api_key=api_key)
                                        uploaded_file_ocr = ocr_client.files.upload(file=temp_ocr_path)
                                        while uploaded_file_ocr.state.name == "PROCESSING":
                                            time.sleep(1)
                                            uploaded_file_ocr = ocr_client.files.get(name=uploaded_file_ocr.name)
                                        ocr_prompt = f"Realiza el reconocimiento de texto (OCR) de la página {preview_page_num} de este documento y devuelve únicamente el texto transcrito en español."
                                        response_ocr = ocr_client.models.generate_content(
                                            model="gemini-3.5-flash",
                                            contents=[uploaded_file_ocr, ocr_prompt]
                                        )
                                        st.session_state.ocr_preview_cache[cache_key] = response_ocr.text
                                        try:
                                            ocr_client.files.delete(name=uploaded_file_ocr.name)
                                        except Exception:
                                            pass
                                    else:
                                        ocr_client = genai.Client(api_key=api_key)
                                        uploaded_file_ocr = ocr_client.files.upload(file=temp_ocr_path)
                                        while uploaded_file_ocr.state.name == "PROCESSING":
                                            time.sleep(1)
                                            uploaded_file_ocr = ocr_client.files.get(name=uploaded_file_ocr.name)
                                        ocr_prompt = f"Realiza el reconocimiento de texto (OCR) de la página {preview_page_num} de este documento y devuelve únicamente el texto transcrito en español."
                                        response_ocr = ocr_client.models.generate_content(
                                            model=selected_model,
                                            contents=[uploaded_file_ocr, ocr_prompt]
                                        )
                                        st.session_state.ocr_preview_cache[cache_key] = response_ocr.text
                                        try:
                                            ocr_client.files.delete(name=uploaded_file_ocr.name)
                                        except Exception:
                                            pass
                                            
                                    if os.path.exists(temp_ocr_path):
                                        os.remove(temp_ocr_path)
                                    st.rerun()
                                except Exception as e:
                                    st.error(f"Error en OCR de página: {str(e)}")
                elif cache_key in st.session_state.ocr_preview_cache:
                    st.success("✅ Texto recuperado por OCR:")
                    st.text_area("Contenido extraído", value=st.session_state.ocr_preview_cache[cache_key], height=200)
                else:
                    st.text_area("Contenido extraído de la página", value=selected_page_text, height=200, disabled=True)
                    
            st.markdown("---")
            force_ocr = st.checkbox(
                "Forzar OCR con IA (Recomendado para libros escaneados o PDFs con páginas en blanco)",
                value=False,
                key="force_ocr_cb"
            )
            
            reading_name = st.text_input(
                "Nombre para guardar la lectura",
                value=uploaded_reading_file.name.rsplit(".", 1)[0],
                key="save_reading_name_input"
            )
            
            if st.button("💾 Guardar Lectura en Historial", key="btn_save_reading_history", use_container_width=True):
                # Guardar el archivo físicamente
                temp_doc_path = f"temp_reading_{uploaded_reading_file.name}"
                with open(temp_doc_path, "wb") as f:
                    f.write(uploaded_reading_file.getbuffer())
                    
                loader = show_full_screen_loader("Analizando y guardando lectura de apoyo...")
                try:
                    # Cargar páginas
                    pages = []
                    suffix = Path(uploaded_reading_file.name).suffix.lower()
                    
                    if force_ocr:
                        # OCR Masivo utilizando Gemini API
                        if not api_key:
                            st.error("🔑 El OCR masivo requiere configurar tu Gemini API Key en el panel lateral.")
                        else:
                            update_full_screen_loader(loader, "⚠️ Iniciando el Reconocimiento de Texto (OCR) con la Inteligencia Artificial de Gemini...")
                            client_ocr = genai.Client(api_key=api_key)
                            uploaded_file_ocr = client_ocr.files.upload(file=temp_doc_path)
                            while uploaded_file_ocr.state.name == "PROCESSING":
                                update_full_screen_loader(loader, "Procesando documento en los servidores de Gemini para OCR...")
                                time.sleep(3)
                                uploaded_file_ocr = client_ocr.files.get(name=uploaded_file_ocr.name)
                                
                            total_pages = len(extract_reading_pages(temp_doc_path, suffix))
                            window_size = 15
                            import math
                            num_chunks = math.ceil(total_pages / window_size)
                            
                            for chunk_idx in range(num_chunks):
                                start_p = chunk_idx * window_size + 1
                                end_p = min((chunk_idx + 1) * window_size, total_pages)
                                
                                update_full_screen_loader(loader, f"Transcribiendo páginas {start_p} a {end_p} de {total_pages} con OCR por IA...")
                                prompt_ocr = f"""
Te proporciono un documento PDF escaneado. Por favor, lee las páginas de la {start_p} a la {end_p} de este documento y realiza una transcripción literal y completa del texto en español.
Devuelve el resultado estrictamente en formato JSON como una lista de objetos. Cada objeto debe tener la estructura exacta:
{{
  "page_num": <número de página entero>,
  "content": "<texto transcrito de la página>"
}}
No añadas formato de bloque de código Markdown o explicaciones fuera del JSON. Devuelve únicamente el JSON.
"""
                                response_ocr = client_ocr.models.generate_content(
                                    model="gemini-3.5-flash",
                                    contents=[uploaded_file_ocr, prompt_ocr]
                                )
                                response_text_ocr = response_ocr.text.strip()
                                if response_text_ocr.startswith("```json"):
                                    response_text_ocr = response_text_ocr[7:]
                                if response_text_ocr.endswith("```"):
                                    response_text_ocr = response_text_ocr[:-3]
                                response_text_ocr = response_text_ocr.strip()
                                
                                try:
                                    ocr_data = json.loads(response_text_ocr)
                                    for item in ocr_data:
                                        pages.append({
                                            "page_num": int(item["page_num"]),
                                            "content": str(item["content"])
                                        })
                                except Exception as json_err:
                                    st.warning(f"Error al analizar JSON del bloque {start_p}-{end_p}, reintentando de forma local...")
                                    # Fallback local
                                    local_pages = extract_reading_pages(temp_doc_path, suffix)
                                    pages.extend(local_pages[start_p-1:end_p])
                                    
                            try:
                                client_ocr.files.delete(name=uploaded_file_ocr.name)
                            except Exception:
                                pass
                    else:
                        pages = extract_reading_pages(temp_doc_path, suffix)
                        
                    # Guardar archivo JSON local
                    save_reading_data(
                        name=reading_name.strip(),
                        filename=uploaded_reading_file.name,
                        pages=pages
                    )
                    st.success(f"💾 Lectura '{reading_name.strip()}' guardada con éxito en el historial de {active_sub}!")
                    time.sleep(2)
                    st.rerun()
                except Exception as e:
                    st.error(f"❌ Error al guardar la lectura: {str(e)}")
                finally:
                    if os.path.exists(temp_doc_path):
                        os.remove(temp_doc_path)
                    loader.empty()
                    
    with col_saved_readings:
        st.markdown("### 📚 Lecturas Guardadas")
        
        readings_list = list_saved_readings()
        
        if not readings_list:
            st.warning("💡 Aún no tienes lecturas guardadas. Sube una lectura en el panel anterior.")
        else:
            selected_reading = st.selectbox(
                "Selecciona una lectura de apoyo",
                options=readings_list,
                key="reading_select"
            )
            
            reading_data = load_reading_data(selected_reading)
            
            if "reading_summary" not in st.session_state:
                st.session_state.reading_summary = None
            if "reading_easy_text" not in st.session_state:
                st.session_state.reading_easy_text = None
            if "extracted_chapter" not in st.session_state:
                st.session_state.extracted_chapter = None
            if "reading_last_action" not in st.session_state:
                st.session_state.reading_last_action = None
            
            reading_action = st.radio(
                "Elige la herramienta a utilizar:",
                options=[
                    "📝 Generar Resumen de la Lectura",
                    "🔎 Modo Lectura Fácil (Legibilidad)",
                    "✂️ Extraer Capítulos / Secciones",
                    "✏️ Cambiar Nombre",
                    "🗑️ Eliminar Lectura"
                ],
                horizontal=True,
                key="reading_action_radio"
            )
            
            if st.session_state.reading_last_action != reading_action:
                st.session_state.reading_summary = None
                st.session_state.reading_easy_text = None
                st.session_state.extracted_chapter = None
                st.session_state.reading_last_action = reading_action
            
            # --- ACCIÓN 1: GENERAR RESUMEN ---
            if reading_action == "📝 Generar Resumen de la Lectura":
                st.markdown("##### 📝 Resumen con Inteligencia Artificial")
                
                reading_depth = st.selectbox(
                    "Profundidad del resumen de la lectura",
                    options=[
                        "Completo y Detallado (Estudio exhaustivo)",
                        "Intermedio (Estándar equilibrado)",
                        "Sencillo (Acotado y directo)"
                    ],
                    index=1,
                    key="reading_depth_select"
                )
                
                if st.button("🤖 Generar Resumen", key="btn_gen_read_sum"):
                    if not api_key and "gpt-" not in selected_model:
                        st.error("🔑 Ingresa tu Gemini API Key en el panel lateral.")
                    elif "gpt-" in selected_model and not openai_api_key:
                        st.error("🔑 Ingresa tu OpenAI API Key en el panel lateral.")
                    else:
                        loader = show_full_screen_loader("Generando resumen de la lectura con Inteligencia Artificial...")
                        try:
                            full_reading_text = "\n\n".join([f"--- PÁGINA {p['page_num']} ---\n{p['content']}" for p in reading_data["pages"]])
                            
                            if "Completo" in reading_depth:
                                instructions = "Genera un resumen de la lectura con alta profundidad. Explaya cada concepto con al menos dos o tres párrafos detallados e incluye de forma textual los ejemplos o casos prácticos expuestos en el texto."
                            elif "Intermedio" in reading_depth:
                                instructions = "Genera un resumen estructurado, equilibrado y claro, sintetizando los capítulos y temas principales de la lectura."
                            else:
                                instructions = "Genera un resumen conciso y directo al grano, utilizando listas con viñetas cortas y párrafos breves sobre las ideas esenciales."
                                
                            prompt = f"""
Actúa como un profesor universitario. Te proporciono el texto de un material de lectura de apoyo.
Tu tarea es generar un resumen académico de esta lectura basado en las siguientes directrices:
- Nivel de detalle: {reading_depth}
- Instrucciones: {instructions}

Estructura el resumen con un título descriptivo, una introducción corta, el desglose ordenado de temas/capítulos principales, y un glosario final si hay términos técnicos relevantes.

### TEXTO COMPLETO DE LA LECTURA:
{full_reading_text}
"""
                            if "gpt-" in selected_model:
                                openai_client = OpenAI(api_key=openai_api_key)
                                response = openai_client.chat.completions.create(
                                    model=selected_model,
                                    messages=[{"role": "user", "content": prompt}]
                                )
                                response_text = response.choices[0].message.content
                            else:
                                client = genai.Client(api_key=api_key)
                                response = client.models.generate_content(
                                    model=selected_model,
                                    contents=[prompt]
                                )
                                response_text = response.text
                                
                            st.session_state.reading_summary = response_text
                        except Exception as e:
                            st.error(handle_api_error(e))
                        finally:
                            loader.empty()
                            
                if st.session_state.reading_summary:
                    st.markdown("---")
                    st.markdown("##### 📝 Resumen Generado:")
                    st.markdown(st.session_state.reading_summary)
                    
                    st.markdown("#### 📥 Descargar Resumen de Lectura")
                    reading_sum_filename = st.text_input(
                        "Nombre de archivo de descarga del resumen",
                        value=f"resumen_{selected_reading.lower().replace(' ', '_')}",
                        key="reading_sum_filename_input"
                    )
                    
                    pdf_bytes_sum = None
                    try:
                        pdf_bytes_sum = generate_pdf_from_markdown(st.session_state.reading_summary)
                    except Exception as e:
                        st.error(f"Error al generar PDF: {str(e)}")
                        
                    col_sum_md, col_sum_pdf = st.columns(2)
                    with col_sum_md:
                        st.download_button(
                            label="📥 Descargar Resumen (Markdown)",
                            data=st.session_state.reading_summary,
                            file_name=f"{reading_sum_filename}.md",
                            mime="text/markdown",
                            key="read_sum_md_btn",
                            use_container_width=True
                        )
                    with col_sum_pdf:
                        if pdf_bytes_sum:
                            st.download_button(
                                label="📥 Descargar Resumen (PDF)",
                                data=pdf_bytes_sum,
                                file_name=f"{reading_sum_filename}.pdf",
                                mime="application/pdf",
                                key="read_sum_pdf_btn",
                                use_container_width=True
                            )
                            
            # --- ACCIÓN 2: MODO LECTURA FÁCIL ---
            elif reading_action == "🔎 Modo Lectura Fácil (Legibilidad)":
                st.markdown("##### 🔎 Modo Lectura Fácil")
                st.markdown("Ajusta la tipografía, tamaño y espaciado para una lectura cómoda en pantalla o genera un PDF de alta legibilidad.")
                
                col_font, col_size, col_spacing, col_align = st.columns(4)
                with col_font:
                    font_choice = st.selectbox("Tipografía", options=["Sans Serif", "Monospace", "Serif", "Dislexia"], index=0, key="easy_font")
                with col_size:
                    size_choice = st.selectbox("Tamaño", options=["Muy Grande", "Mediana", "Pequeña"], index=1, key="easy_size")
                with col_spacing:
                    spacing_choice = st.selectbox("Espaciado", options=["Espaciado (2.0)", "Medio (1.5)", "Estrecho (1.0)"], index=1, key="easy_spacing")
                with col_align:
                    align_choice = st.selectbox("Alineación", options=["Izquierda", "Centrado", "Justificado"], index=2, key="easy_align")
                    
                theme_choice = st.selectbox(
                    "Tema visual para lectura en pantalla",
                    options=["Clásico (Blanco/Negro)", "Crema (Sepia)", "Oscuro", "Alto Contraste (Negro/Amarillo)"],
                    index=0,
                    key="easy_theme"
                )
                
                full_reading_text = "\n\n".join([f"--- PÁGINA {p['page_num']} ---\n{p['content']}" for p in reading_data["pages"]])
                
                font_style = "Arial, sans-serif"
                if font_choice == "Serif":
                    font_style = "Georgia, serif"
                elif font_choice == "Monospace":
                    font_style = "Courier New, monospace"
                elif font_choice == "Dislexia":
                    font_style = "Comic Sans MS, cursive"
                    
                size_style = "18px"
                if size_choice == "Muy Grande":
                    size_style = "24px"
                elif size_choice == "Pequeña":
                    size_style = "14px"
                    
                spacing_style = "1.5"
                if "2.0" in spacing_choice:
                    spacing_style = "2.0"
                elif "1.0" in spacing_choice:
                    spacing_style = "1.0"
                    
                align_style = "justify"
                if align_choice == "Izquierda":
                    align_style = "left"
                elif align_choice == "Centrado":
                    align_style = "center"
                    
                bg_color = "#ffffff"
                text_color = "#000000"
                if theme_choice == "Crema (Sepia)":
                    bg_color = "#f4ecd8"
                    text_color = "#5b4636"
                elif theme_choice == "Oscuro":
                    bg_color = "#1e1e1e"
                    text_color = "#e0e0e0"
                elif theme_choice == "Alto Contraste (Negro/Amarillo)":
                    bg_color = "#000000"
                    text_color = "#ffff00"
                    
                st.markdown(f"""
                <div style="font-family: {font_style}; font-size: {size_style}; line-height: {spacing_style}; text-align: {align_style}; background-color: {bg_color}; color: {text_color}; padding: 25px; border-radius: 12px; border: 1px solid #ddd; max-height: 500px; overflow-y: auto;">
                {full_reading_text.replace(chr(10), '<br>')}
                </div>
                """, unsafe_allow_html=True)
                
                pdf_easy_bytes = None
                try:
                    pdf_size = 11.5
                    if size_choice == "Muy Grande":
                        pdf_size = 14.0
                    elif size_choice == "Pequeña":
                        pdf_size = 9.5
                        
                    pdf_spacing = 1.5
                    if "2.0" in spacing_choice:
                        pdf_spacing = 2.0
                    elif "1.0" in spacing_choice:
                        pdf_spacing = 1.0
                        
                    pdf_easy_bytes = generate_easy_read_pdf(
                        text=full_reading_text,
                        font_family=font_choice,
                        font_size=pdf_size,
                        line_spacing=pdf_spacing,
                        alignment=align_choice
                    )
                except Exception as e:
                    st.error(f"Error al generar PDF de Lectura Fácil: {str(e)}")
                    
                if pdf_easy_bytes:
                    st.download_button(
                        label="📥 Descargar PDF Lectura Fácil",
                        data=pdf_easy_bytes,
                        file_name=f"lectura_facil_{selected_reading.lower().replace(' ', '_')}.pdf",
                        mime="application/pdf",
                        key="easy_pdf_btn",
                        use_container_width=True
                    )
                    
            # --- ACCIÓN 3: EXTRAER CAPÍTULOS / SECCIONES ---
            elif reading_action == "✂️ Extraer Capítulos / Secciones":
                st.markdown("##### ✂️ Extraer Sección o Capítulo")
                
                extract_type = st.radio(
                    "Método de extracción:",
                    options=["Por Rango de Páginas (Local)", "Por Nombre de Capítulo/Tema (IA)"],
                    horizontal=True,
                    key="extract_type_select"
                )
                
                if extract_type == "Por Rango de Páginas (Local)":
                    col_p_start, col_p_end = st.columns(2)
                    with col_p_start:
                        p_start = st.number_input("Página de inicio", min_value=1, max_value=len(reading_data["pages"]), value=1)
                    with col_p_end:
                        p_end = st.number_input("Página de fin", min_value=1, max_value=len(reading_data["pages"]), value=min(10, len(reading_data["pages"])))
                        
                    if st.button("✂️ Extraer Rango", key="btn_extract_range"):
                        if p_start > p_end:
                            st.error("❌ La página de inicio no puede ser mayor que la de fin.")
                        else:
                            extracted_pages = [p["content"] for p in reading_data["pages"] if p_start <= p["page_num"] <= p_end]
                            st.session_state.extracted_chapter = "\n\n".join(extracted_pages)
                            st.success(f"✂️ Páginas {p_start} a {p_end} extraídas correctamente.")
                else:
                    chapter_query = st.text_input(
                        "Nombre o tema del capítulo a extraer",
                        help="Escribe el nombre del capítulo o el tema principal que deseas extraer. La IA buscará en qué páginas se encuentra y lo transcribirá íntegramente."
                    )
                    
                    scan_window = st.slider(
                        "Tamaño de la ventana de páginas a enviar a la IA:",
                        min_value=15,
                        max_value=200,
                        value=50,
                        help="Define cuántas páginas a la redonda del inicio del capítulo se enviarán a Gemini. Si el capítulo es muy extenso, incrementa este valor."
                    )
                    
                    if st.button("✂️ Extraer Capítulo con IA", key="btn_extract_chapter_ia"):
                        if not chapter_query.strip():
                            st.warning("⚠️ Ingresa el nombre del capítulo.")
                        elif not api_key and "gpt-" not in selected_model:
                            st.error("🔑 Ingresa tu Gemini API Key en el panel lateral.")
                        elif "gpt-" in selected_model and not openai_api_key:
                            st.error("🔑 Ingresa tu OpenAI API Key en el panel lateral.")
                        else:
                            loader = show_full_screen_loader(f"Localizando y extrayendo sección '{chapter_query}' con IA...")
                            try:
                                full_reading_text = get_relevant_pages_for_extraction(reading_data, chapter_query, window_size=scan_window)
                                
                                prompt = f"""
Actúa como un asistente editorial académico de alta precisión. Te proporciono el texto de las páginas de un libro.
Tu tarea consiste en extraer de forma íntegra, textual y completa únicamente el desarrollo y contenido del capítulo/sección indicado:
- Capítulo/Sección a extraer: "{chapter_query}"

### DIRECTRICES DE EXTRACCIÓN CRÍTICAS:
1. **Ignora los falsos positivos**: No extraigas resúmenes, descripciones cortas, esquemas o menciones del capítulo que aparezcan en el Prólogo, Introducción, o la sección de "Estructura de la obra" o "Contenido del CD". Debes buscar el **inicio real del cuerpo del capítulo** (donde se desarrolla la lectura y teoría en sí).
2. **Extrae el texto de forma completa**: Transcribe textualmente y de forma íntegra los párrafos, subtítulos, ejemplos y explicaciones de ese capítulo. No hagas resúmenes ni recortes, extrae todo el contenido.
3. Devuelve la extracción con un formato limpio en español.

### TEXTO DE LAS PÁGINAS SELECCIONADAS DEL LIBRO:
{full_reading_text}
"""
                                if "gpt-" in selected_model:
                                    openai_client = OpenAI(api_key=openai_api_key)
                                    response = openai_client.chat.completions.create(
                                        model=selected_model,
                                        messages=[{"role": "user", "content": prompt}]
                                    )
                                    response_text = response.choices[0].message.content
                                else:
                                    client = genai.Client(api_key=api_key)
                                    response = client.models.generate_content(
                                        model=selected_model,
                                        contents=[prompt]
                                    )
                                    response_text = response.text
                                    
                                st.session_state.extracted_chapter = response_text
                                st.success(f"✂️ Sección '{chapter_query}' localizada y extraída con éxito.")
                            except Exception as e:
                                st.error(handle_api_error(e))
                            finally:
                                loader.empty()
                                
                if st.session_state.extracted_chapter:
                    st.markdown("---")
                    st.markdown("##### 📍 Contenido Extraído:")
                    st.markdown(st.session_state.extracted_chapter)
                    
                    st.markdown("#### 📥 Descargar Sección Extraída")
                    chapter_filename = st.text_input(
                        "Nombre de archivo de descarga de la sección",
                        value=f"extraccion_{selected_reading.lower().replace(' ', '_')}",
                        key="extracted_chapter_filename_input"
                    )
                    
                    pdf_bytes_chap = None
                    try:
                        pdf_bytes_chap = generate_pdf_from_markdown(st.session_state.extracted_chapter)
                    except Exception as e:
                        st.error(f"Error al generar PDF: {str(e)}")
                        
                    col_chap_md, col_chap_pdf = st.columns(2)
                    with col_chap_md:
                        st.download_button(
                            label="📥 Descargar Sección (Markdown)",
                            data=st.session_state.extracted_chapter,
                            file_name=f"{chapter_filename}.md",
                            mime="text/markdown",
                            key="chap_md_btn",
                            use_container_width=True
                        )
                    with col_chap_pdf:
                        if pdf_bytes_chap:
                            st.download_button(
                                label="📥 Descargar Sección (PDF)",
                                data=pdf_bytes_chap,
                                file_name=f"{chapter_filename}.pdf",
                                mime="application/pdf",
                                key="chap_pdf_btn",
                                use_container_width=True
                            )
                            
            # --- ACCIÓN 4: CAMBIAR NOMBRE ---
            elif reading_action == "✏️ Cambiar Nombre":
                st.markdown("##### ✏️ Cambiar Nombre de la Lectura")
                new_reading_name_input = st.text_input(
                    "Nuevo nombre de la lectura",
                    value=selected_reading,
                    help="Modifica el nombre con el que se guarda esta lectura en tu historial local."
                )
                
                if st.button("✏️ Cambiar Nombre de Lectura", key="btn_rename_reading_action"):
                    if not new_reading_name_input.strip():
                        st.warning("⚠️ Ingresa un nombre válido.")
                    else:
                        try:
                            old_reading_path = get_readings_dir() / f"{selected_reading}.json"
                            new_reading_path = get_readings_dir() / f"{new_reading_name_input.strip()}.json"
                            
                            if new_reading_path.exists():
                                st.error("❌ Ya existe una lectura con ese nombre en esta materia.")
                            else:
                                r_data = load_reading_data(selected_reading)
                                if r_data:
                                    r_data["name"] = new_reading_name_input.strip()
                                    with open(new_reading_path, 'w', encoding='utf-8') as f:
                                        json.dump(r_data, f, ensure_ascii=False, indent=4)
                                    os.unlink(old_reading_path)
                                    st.success(f"✏️ Lectura renombrada con éxito a '{new_reading_name_input.strip()}'")
                                    time.sleep(2)
                                    st.rerun()
                        except Exception as e:
                            st.error(f"❌ Error al renombrar la lectura: {str(e)}")
                            
            # --- ACCIÓN 5: ELIMINAR LECTURA ---
            elif reading_action == "🗑️ Eliminar Lectura":
                st.markdown("##### 🗑️ Eliminar Lectura de la Base de Datos")
                st.write(f"¿Estás seguro de que deseas eliminar permanentemente la lectura `{selected_reading}`?")
                st.warning("⚠️ Esta acción eliminará el texto y el registro de la lectura de forma permanente.")
                
                if st.button("🗑️ Eliminar lectura de forma definitiva", key="btn_del_read_permanent"):
                    try:
                        file_path = get_readings_dir() / f"{selected_reading}.json"
                        if file_path.exists():
                            os.unlink(file_path)
                            st.success(f"🗑️ Lectura '{selected_reading}' eliminada con éxito.")
                            time.sleep(2)
                            st.rerun()
                    except Exception as e:
                        st.error(f"❌ Error al eliminar la lectura: {str(e)}")

# =====================================================================
# PESTAÑA 5: TUTOR IA (CHAT CON CONTEXTO CON SOPORTE OPENAI Y GEMINI)
# =====================================================================
with tab_chat:
    active_sub = get_active_subject()
    
    col_chat_title, col_chat_clear = st.columns([5, 1])
    with col_chat_title:
        st.markdown(f"### 💬 Tutor Académico Personal - Materia: **{active_sub}**")
    with col_chat_clear:
        chat_key = f"chat_history_{active_sub}"
        if chat_key in st.session_state and st.session_state[chat_key]:
            if st.button("🗑️ Limpiar Chat", key="btn_clear_chat"):
                st.session_state[chat_key] = []
                st.rerun()

    compiled_context = []
    classes = list_saved_classes()
    for cls in classes:
        c_data = load_class_data(cls)
        if c_data:
            compiled_context.append(f"### RESUMEN DE CLASE: {c_data['name']}\n{c_data['summary']}")
    
    readings = list_saved_readings()
    for rd in readings:
        r_data = load_reading_data(rd)
        if r_data:
            r_text = get_reading_consolidated_text(r_data, max_chars=40000)
            compiled_context.append(f"### LECTURA DE APOYO: {r_data['name']}\n{r_text}")

    # Verificar presencia de documentos en Supabase DB
    supabase_doc_count = 0
    try:
        from supabase_client import get_supabase_client
        client = get_supabase_client()
        res_db = client.table("documents").select("id").ilike("title", f"%[{active_sub}]%").execute()
        if res_db.data:
            supabase_doc_count = len(res_db.data)
    except Exception:
        supabase_doc_count = 0

    if not classes and not readings and supabase_doc_count == 0:
        st.warning("💡 Esta materia no tiene clases ni lecturas guardadas todavía. Puedes conversar con el Tutor IA, pero sus respuestas serán generales ya que no hay material específico de estudio.")
        subject_context_str = "No hay materiales de estudio específicos cargados."
    else:
        subject_context_str = "\n\n=========================================\n\n".join(compiled_context)
        if supabase_doc_count > 0:
            st.info(f"⚡ Conectado a Supabase BD ({supabase_doc_count} documentos y vectores RAG activos para **{active_sub}**).")
        else:
            st.info(f"📚 Conectado al material de estudio local de **{active_sub}** ({len(classes)} clases, {len(readings)} lecturas de apoyo).")

    if chat_key not in st.session_state:
        st.session_state[chat_key] = []

    chat_container = st.container()
    with chat_container:
        for message in st.session_state[chat_key]:
            with st.chat_message(message["role"]):
                st.markdown(message["content"])

    if prompt := st.chat_input("Escribe tu duda académica sobre esta materia..."):
        st.session_state[chat_key].append({"role": "user", "content": prompt})
        with chat_container:
            with st.chat_message("user"):
                st.markdown(prompt)

        is_openai = "gpt-" in selected_model
        if is_openai and not openai_api_key:
            st.error("🔑 Ingresa una OpenAI API Key en el panel lateral.")
        elif not is_openai and not api_key:
            st.error("🔑 Ingresa una Gemini API Key en el panel lateral.")
        else:
            # Intentar recuperar fragmentos RAG por similitud vectorial desde Supabase DB
            rag_snippets = []
            try:
                from rag_service import search_similar_chunks
                from supabase_client import get_supabase_client
                client = get_supabase_client()
                res_docs = client.table("documents").select("id, title").ilike("title", f"%[{active_sub}]%").execute()
                if res_docs.data:
                    for d in res_docs.data:
                        matched = search_similar_chunks(d["id"], prompt, top_k=3, api_key=api_key)
                        for m in matched:
                            sim = m.get("similarity", 0)
                            content = m.get("content", "").strip()
                            if content:
                                rag_snippets.append(f"📌 [Fuente: {d['title']} | Similitud: {sim:.2f}]:\n{content}")
            except Exception as e:
                print(f"[RAG Query Notice] Fallback: {e}")
                
            rag_str = "\n\n".join(rag_snippets) if rag_snippets else ""
            if rag_str:
                full_context = f"### FRAGMENTOS RECUPERADOS DESDE SUPABASE RAG BD:\n{rag_str}\n\n### MATERIAL DE APOYO ADICIONAL:\n{subject_context_str}"
            else:
                full_context = subject_context_str

            system_instruction = f"""
Actúa como un tutor académico de alta competencia en el ámbito universitario. Tu rol es resolver de forma didáctica, clara y rigurosa cualquier duda sobre el material de estudio de la materia activa.
Tienes acceso al contenido consolidado de la materia activa y fragmentos vectoriales recuperados de Supabase RAG que se te adjuntan a continuación. Utilízalo como tu fuente principal de verdad académica.

### CONTEXTO DE LA MATERIA ACTIVA:
{full_context}

### INSTRUCCIONES:
1. Responde de forma precisa y enfocándote en los conceptos del material.
2. Si la pregunta no se relaciona con la materia o no puedes responderla con el material, indícalo de forma honesta, pero ofrece guiar al estudiante sobre cómo buscar en el material existente.
3. Mantén un tono formal, motivador y sumamente educativo.
"""
            
            if is_openai:
                openai_client = OpenAI(api_key=openai_api_key)
                messages = [{"role": "system", "content": system_instruction}]
                for msg in st.session_state[chat_key]:
                    messages.append({"role": msg["role"], "content": msg["content"]})
                    
                with st.spinner("Tutor IA escribiendo..."):
                    try:
                        response = openai_client.chat.completions.create(
                            model=selected_model,
                            messages=messages
                        )
                        assistant_response = response.choices[0].message.content
                        st.session_state[chat_key].append({"role": "assistant", "content": assistant_response})
                        st.rerun()
                    except Exception as e:
                        st.error(f"❌ Error al consultar a OpenAI: {str(e)}")
            else:
                client = genai.Client(api_key=api_key)
                contents = []
                for message in st.session_state[chat_key][:-1]:
                    role = "user" if message["role"] == "user" else "model"
                    contents.append(
                        genai.types.Content(
                            role=role,
                            parts=[genai.types.Part.from_text(text=message["content"])]
                        )
                    )
                contents.append(
                    genai.types.Content(
                        role="user",
                        parts=[genai.types.Part.from_text(text=prompt)]
                    )
                )
                
                with st.spinner("Tutor IA escribiendo..."):
                    try:
                        response = client.models.generate_content(
                            model=selected_model,
                            contents=contents,
                            config=genai.types.GenerateContentConfig(
                                system_instruction=system_instruction
                            )
                        )
                        assistant_response = response.text
                        st.session_state[chat_key].append({"role": "assistant", "content": assistant_response})
                        st.rerun()
                    except Exception as e:
                        st.error(handle_api_error(e))
